from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x0A, 0x1A)   # near-black purple
CARD      = RGBColor(0x1E, 0x15, 0x35)   # card bg
VIOLET    = RGBColor(0x7C, 0x3A, 0xED)   # accent 1
CYAN      = RGBColor(0x06, 0xB6, 0xD4)   # accent 2
ORANGE    = RGBColor(0xF9, 0x73, 0x16)   # accent 3
PINK      = RGBColor(0xEC, 0x48, 0x99)   # accent 4
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # accent 5
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC4, 0xB5, 0xFD)   # light purple text
MUTED     = RGBColor(0x9D, 0x8E, 0xC2)   # muted text
ROSE      = RGBColor(0xF4, 0x3F, 0x5E)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return sl

def box(sl, x, y, w, h, color=None, alpha=None):
    shape = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    return shape

def txt(sl, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def multiline_txt(sl, lines, x, y, w, h, size=14, color=WHITE,
                  bold=False, align=PP_ALIGN.LEFT, line_space=1.2):
    """lines = list of (text, size, bold, color) or just strings"""
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, str):
            ls, lb, lc = size, bold, color
            lt = line
        else:
            lt, ls, lb, lc = line
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = lt
        run.font.size = Pt(ls)
        run.font.bold = lb
        run.font.color.rgb = lc
    return txb

def accent_bar(sl, x, y, w=0.06, h=0.5, color=VIOLET):
    """Thin vertical accent bar"""
    box(sl, x, y, w, h, color)

def section_label(sl, text, x=0.5, y=0.28):
    txt(sl, text, x, y, 5, 0.3, size=9, bold=True, color=VIOLET)

def slide_headline(sl, text, x=0.5, y=0.6, w=12.3, size=34):
    txt(sl, text, x, y, w, 1.0, size=size, bold=True, color=WHITE)

def divider_line(sl, x, y, w, color=VIOLET):
    from pptx.util import Pt as Pt2
    ln = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def card_rect(sl, x, y, w, h, color=CARD, border=VIOLET):
    shape = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    return shape

def stat_card(sl, value, label, x, y, accent=VIOLET):
    card_rect(sl, x, y, 2.8, 1.5, CARD, accent)
    txt(sl, value, x+0.15, y+0.12, 2.5, 0.75, size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txt(sl, label, x+0.15, y+0.85, 2.5, 0.5, size=10, bold=False, color=MUTED, align=PP_ALIGN.CENTER)

def tag(sl, text, x, y, color=VIOLET):
    card_rect(sl, x, y, len(text)*0.11+0.2, 0.3, color, color)
    txt(sl, text, x+0.1, y+0.04, len(text)*0.11+0.1, 0.25, size=9, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()

# Mosaic title — big gradient simulation with two overlapping layers
txt(sl, "MOSAIC", 1.5, 1.6, 10, 2.2, size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Underlay in violet for depth
txt(sl, "MOSAIC", 1.52, 1.62, 10, 2.2, size=96, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
txt(sl, "MOSAIC", 1.5, 1.6, 10, 2.2, size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Ground truth at scale.", 1.5, 3.65, 10.3, 0.6, size=22, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)
txt(sl, "Physical world intelligence — on demand, AI-extracted, instantly.", 1.5, 4.2, 10.3, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)

divider_line(sl, 4.0, 4.85, 5.3, VIOLET)

txt(sl, "[Founder Name]  ·  [Founder Name]", 1.5, 5.1, 10.3, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "hello@mosaic-intel.com  |  mosaic-intel.com  |  Confidential · March 2026", 1.5, 5.5, 10.3, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Decorative dots (mosaic tile illusion)
for i in range(6):
    for j in range(3):
        c = [VIOLET, CYAN, PINK][j % 3]
        b = box(sl, 0.3 + i*0.18, 0.2 + j*0.18, 0.09, 0.09, c)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=ROSE)
section_label(sl, "THE PROBLEM", 0.55, 0.28)
slide_headline(sl, "The physical world is opaque —\nand everyone pays for it", 0.5, 0.6, size=30)

# 3 problem columns
cols = [
    ("CPG BRANDS", "Lost 4–8% of revenue to\nout-of-stocks they never\nknew about", ROSE),
    ("RETAILERS", "Can't verify compliance\nwithout 2-week delay\nand $200/visit", ORANGE),
    ("CONSUMERS", "Can't trust outdated\nphotos, surveys, and\nstale shelf data", PINK),
]
for i, (title, body, color) in enumerate(cols):
    cx = 0.5 + i * 4.25
    card_rect(sl, cx, 2.05, 4.0, 1.75, CARD, color)
    txt(sl, title, cx+0.2, 2.15, 3.6, 0.4, size=13, bold=True, color=color)
    txt(sl, body, cx+0.2, 2.55, 3.6, 1.1, size=12, color=LIGHT)

# Core problem box
card_rect(sl, 0.5, 4.0, 12.33, 2.15, RGBColor(0x28, 0x18, 0x48), VIOLET)
txt(sl, "THE CORE PROBLEM", 0.8, 4.1, 5, 0.35, size=11, bold=True, color=VIOLET)
problem_lines = [
    "The physical world generates zero data unless someone deliberately captures it.",
    "Field checks cost $50–$300 per visit · Take 2–3 weeks · Cover <5% of locations · Deliver no machine-readable output.",
    "$35B is spent on physical world intelligence per year. Most of it produces data that's already outdated when it arrives.",
]
for i, line in enumerate(problem_lines):
    txt(sl, line, 0.8, 4.5 + i*0.5, 11.8, 0.45, size=12, color=WHITE if i==0 else LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=CYAN)
section_label(sl, "THE SOLUTION", 0.55, 0.28)
slide_headline(sl, "Mosaic turns any smartphone into a\nverified data collection device", 0.5, 0.6, size=28)

steps = [
    ("🏢", "ENTERPRISE POSTS", "Check oat milk shelf at 50\nLondon Tesco stores", VIOLET),
    ("📡", "MOSAIC ROUTES", "12 agents already near those\nstores on their commutes", CYAN),
    ("📸", "AGENTS CAPTURE", "On-device face redaction +\nquality check before upload", GREEN),
    ("🤖", "AI EXTRACTS", "JSON data, compliance scores,\nanomaly flags — in 4 hours", ORANGE),
    ("📊", "INTELLIGENCE DELIVERED", "6 OOS · 3 wrong placement\n41 compliant — actionable", PINK),
]

for i, (icon, title, body, color) in enumerate(steps):
    cx = 0.35 + i * 2.6
    card_rect(sl, cx, 2.1, 2.4, 2.8, CARD, color)
    txt(sl, icon, cx+0.85, 2.2, 0.7, 0.5, size=20, align=PP_ALIGN.CENTER)
    txt(sl, title, cx+0.1, 2.75, 2.2, 0.45, size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(sl, body, cx+0.1, 3.2, 2.2, 1.0, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "›", cx+2.4, 3.2, 0.2, 0.5, size=22, color=VIOLET, align=PP_ALIGN.CENTER)

# Bottom tagline
divider_line(sl, 0.5, 5.15, 12.33, VIOLET)
txt(sl, "Consumer posts in 90 seconds.  Enterprise integrates via API.  Everyone gets ground truth — faster and cheaper than anything else.",
    0.5, 5.3, 12.33, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Magic numbers
magic = [("90 sec", "to post a task"), ("4 hours", "to deliver results"), ("45 sec", "per agent task"), ("55%+", "gross margin")]
for i, (num, lbl) in enumerate(magic):
    txt(sl, num, 0.5 + i*3.3, 5.85, 3.0, 0.5, size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    txt(sl, lbl, 0.5 + i*3.3, 6.35, 3.0, 0.3, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — WHY NOW
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=ORANGE)
section_label(sl, "WHY NOW", 0.55, 0.28)
slide_headline(sl, "Three forces converged in 2026 —\nmaking Mosaic possible and defensible", 0.5, 0.6, size=28)

pillars = [
    ("01", "ON-DEVICE AI", "MediaPipe + TFLite run face redaction in real-time on any 2020+ smartphone. GDPR compliance without cloud roundtrip. This wasn't possible before 2024.", VIOLET),
    ("02", "GIG ECONOMY MATURITY", "50M+ gig workers globally. Deliveroo, Uber, DoorDash trained the supply side. We harvest the agent infrastructure they built.", CYAN),
    ("03", "AI EXTRACTION QUALITY", "GPT-4V class models extract structured data from photos with >90% accuracy. A shelf photo becomes JSON data for $0.20 of compute.", GREEN),
]
for i, (num, title, body, color) in enumerate(pillars):
    cx = 0.4 + i * 4.3
    card_rect(sl, cx, 2.05, 4.0, 3.2, CARD, color)
    txt(sl, num, cx+0.2, 2.15, 0.8, 0.7, size=36, bold=True, color=color)
    txt(sl, title, cx+0.2, 2.9, 3.6, 0.45, size=12, bold=True, color=WHITE)
    txt(sl, body, cx+0.2, 3.4, 3.6, 1.7, size=11, color=LIGHT)

card_rect(sl, 0.4, 5.45, 12.53, 1.3, RGBColor(0x1A, 0x10, 0x30), ORANGE)
txt(sl, "These three forces create a 24-month window before the major US platforms notice Europe.",
    0.7, 5.55, 9.0, 0.45, size=14, bold=True, color=WHITE)
txt(sl, "The European competitor (BeMyEye) has no consumer tier. The enterprise AI buyers are ready — they signed AI tools faster in 2025 than any year prior.",
    0.7, 6.0, 11.8, 0.55, size=11, color=LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — PRODUCT
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=GREEN)
section_label(sl, "PRODUCT", 0.55, 0.28)
slide_headline(sl, "The product agents love to use\nand customers can't stop checking", 0.5, 0.6, size=28)

# Customer side
card_rect(sl, 0.4, 2.05, 5.9, 4.6, CARD, VIOLET)
txt(sl, "CUSTOMER", 0.6, 2.15, 5.5, 0.35, size=11, bold=True, color=VIOLET)
txt(sl, "Post a task in 90 seconds", 0.6, 2.55, 5.5, 0.45, size=16, bold=True, color=WHITE)
cust_steps = ["1. Enter location + task type", "2. Select 24hr or 4hr delivery", "3. Pay by credit card", "4. Receive AI-extracted JSON report"]
for i, s in enumerate(cust_steps):
    txt(sl, s, 0.7, 3.1 + i*0.52, 5.3, 0.45, size=12, color=LIGHT)
txt(sl, "Task types: Shelf audit · Crowd check · Stock verification\nBrand compliance · Competitive pricing · Site inspection",
    0.7, 5.3, 5.3, 0.7, size=10, color=MUTED)

# Arrow
txt(sl, "⟷", 6.35, 3.8, 0.65, 0.6, size=28, color=VIOLET, align=PP_ALIGN.CENTER)

# Agent side
card_rect(sl, 7.05, 2.05, 5.9, 4.6, CARD, CYAN)
txt(sl, "AGENT", 7.25, 2.15, 5.5, 0.35, size=11, bold=True, color=CYAN)
txt(sl, "Earn £14.50 on your morning commute", 7.25, 2.55, 5.5, 0.45, size=14, bold=True, color=WHITE)
agent_steps = ["1. Turn on Commute Mode before leaving", "2. Pre-accept 3 tasks along your route", "3. Complete each in 45–90 seconds", "4. Earn £14.50 before arriving at work"]
for i, s in enumerate(agent_steps):
    txt(sl, s, 7.35, 3.1 + i*0.52, 5.3, 0.45, size=12, color=LIGHT)
txt(sl, "Commute Mode: tasks matched to route, not just radius\nPre-acceptance: commit before you leave, earn on the way",
    7.35, 5.3, 5.3, 0.7, size=10, color=MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET SIZE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=CYAN)
section_label(sl, "MARKET SIZE", 0.55, 0.28)
slide_headline(sl, "$35B market · 22% CAGR\nThe largest segment has never had a self-serve product", 0.5, 0.6, size=26)

# TAM/SAM/SOM circles (represented as cards)
circles = [
    ("TAM", "$35B", "Physical world intelligence\nRetail + Mystery shopping +\nInsurance field + AI training data", VIOLET, 0.4, 2.0, 4.0, 3.8),
    ("SAM", "$8.5B", "Crowdsourced, on-demand,\nvisual intelligence segment\nSmartphone-first platform", CYAN, 4.7, 2.3, 3.5, 3.3),
    ("SOM (Yr 3)", "$350M", "50K enterprise + 2M SMB\n+ 50M consumer credits", GREEN, 8.45, 2.6, 3.0, 2.8),
]
for label, val, desc, color, cx, cy, cw, ch in circles:
    card_rect(sl, cx, cy, cw, ch, CARD, color)
    txt(sl, label, cx+0.2, cy+0.15, cw-0.3, 0.35, size=11, bold=True, color=color)
    txt(sl, val, cx+0.2, cy+0.55, cw-0.3, 0.75, size=30, bold=True, color=WHITE)
    txt(sl, desc, cx+0.2, cy+1.35, cw-0.3, ch-1.5, size=10, color=MUTED)

# Hidden upside
card_rect(sl, 0.4, 5.95, 12.53, 1.2, RGBColor(0x1A, 0x10, 0x30), ORANGE)
txt(sl, "THE CATEGORY THAT DOESN'T EXIST YET:", 0.7, 6.05, 5.0, 0.35, size=10, bold=True, color=ORANGE)
txt(sl, "BeMyEye (enterprise-only): $35M in 15 years.   Field Agent (enterprise-only): no consumer tier.   Trax: hardware only.",
    0.7, 6.42, 11.8, 0.4, size=11, color=LIGHT)
txt(sl, "The consumer/SMB tier has never been served. That's $3–5B of market that doesn't yet exist as a category.",
    0.7, 6.85, 11.8, 0.25, size=11, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=PINK)
section_label(sl, "BUSINESS MODEL", 0.55, 0.28)
slide_headline(sl, "Three tiers. One platform. 55%+ gross margin.", 0.5, 0.6, size=30)

tiers = [
    ("TIER 1", "CONSUMER", "$5–$25 / task", "Credit wallet\nNo commitment\nApp store / viral", "52%\ngross margin", VIOLET),
    ("TIER 2", "SMB SELF-SERVE", "$49–$749 / mo", "Monthly subscription\nCredit card only\nNo sales call ever", "62%\ngross margin", CYAN),
    ("TIER 3", "ENTERPRISE", "$2.5K–$50K+ / mo", "Annual contracts\nCustom SLAs\nFounder → AE sales", "58%\ngross margin", ORANGE),
]
for i, (tier, name, price, model, margin, color) in enumerate(tiers):
    cx = 0.4 + i * 4.3
    card_rect(sl, cx, 1.95, 4.0, 3.7, CARD, color)
    txt(sl, tier, cx+0.2, 2.05, 3.6, 0.3, size=10, bold=True, color=color)
    txt(sl, name, cx+0.2, 2.4, 3.6, 0.5, size=18, bold=True, color=WHITE)
    txt(sl, price, cx+0.2, 2.95, 3.6, 0.45, size=14, bold=True, color=color)
    txt(sl, model, cx+0.2, 3.5, 3.6, 0.9, size=11, color=LIGHT)
    card_rect(sl, cx+0.2, 4.55, 3.6, 0.8, RGBColor(0x12, 0x0A, 0x20), color)
    txt(sl, margin, cx+0.2, 4.6, 3.6, 0.7, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

# Unit economics strip
card_rect(sl, 0.4, 5.85, 7.8, 1.3, RGBColor(0x1A, 0x10, 0x30), VIOLET)
txt(sl, "UNIT ECONOMICS (BLENDED):", 0.65, 5.95, 4.0, 0.3, size=10, bold=True, color=VIOLET)
econ = "Revenue $14  ·  Agent payout (40%) –$5.60  ·  AI processing –$0.20  ·  Infra –$0.36  →  Gross profit $7.14  (51%)"
txt(sl, econ, 0.65, 6.35, 7.3, 0.55, size=11, color=LIGHT)

# Training data callout
card_rect(sl, 8.35, 5.85, 4.6, 1.3, RGBColor(0x1A, 0x10, 0x30), GREEN)
txt(sl, "BONUS REVENUE STREAM", 8.55, 5.95, 4.2, 0.3, size=10, bold=True, color=GREEN)
txt(sl, "AI training data licensing\nMonth 36: $2.4M/mo at 95% margin", 8.55, 6.35, 4.2, 0.7, size=11, color=LIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — TRACTION (PLACEHOLDER)
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=GREEN)
section_label(sl, "TRACTION", 0.55, 0.28)
slide_headline(sl, "Built by 2 people + AI in 8 weeks", 0.5, 0.6, size=30)

# Stat cards
stat_card(sl, "[X] Agents", "Verified in [# cities]", 0.4, 1.9, VIOLET)
stat_card(sl, "[X] Tasks", "Completed to date", 3.45, 1.9, CYAN)
stat_card(sl, "[X] Customers", "Paying + active pilots", 6.5, 1.9, ORANGE)
stat_card(sl, "$[X] MRR", "Across all 3 tiers", 9.55, 1.9, PINK)

# What we built
card_rect(sl, 0.4, 3.65, 12.53, 2.5, CARD, VIOLET)
txt(sl, "WHAT 2 PEOPLE + AI BUILT IN 8 WEEKS   ·   $30K/month burn rate", 0.7, 3.75, 10.0, 0.35, size=11, bold=True, color=VIOLET)
built = [
    ("Agent App (iOS + Android)", "Flutter — task discovery, commute mode, on-device redaction"),
    ("Customer Portal", "Next.js + Supabase — task posting, live map, AI result delivery"),
    ("AI Extraction Pipeline", "GPT-4V + custom validation — structured JSON in <3 min"),
    ("Route-Matching Algorithm", "Geospatial corridor matching — agent pre-acceptance before leaving home"),
]
for i, (name, desc) in enumerate(built):
    col = i % 2
    row = i // 2
    cx = 0.7 + col * 6.2
    cy = 4.2 + row * 0.65
    txt(sl, f"✓ {name}", cx, cy, 3.5, 0.3, size=11, bold=True, color=GREEN)
    txt(sl, desc, cx+0.2, cy+0.3, 5.5, 0.3, size=10, color=MUTED)

txt(sl, "Traditional team equivalent: 10 engineers, $400K burn, 6 months",
    0.7, 5.85, 12.0, 0.35, size=11, bold=False, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — GO TO MARKET
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=ORANGE)
section_label(sl, "GO-TO-MARKET", 0.55, 0.28)
slide_headline(sl, "Beachhead → expansion → the physical world intelligence default", 0.5, 0.6, size=24)

phases = [
    ("PHASE 1\nNOW → MONTH 6", "BEACHHEAD", "CPG shelf audits\n2 US cities + London\n10 enterprise customers\n500 agents / city", VIOLET),
    ("PHASE 2\nMONTH 6 → 18", "MULTI-CITY", "8 US + 5 EU cities\nSMB self-serve launches\n15,000 active agents\nContent + PLG growth", CYAN),
    ("PHASE 3\nMONTH 18+", "CATEGORY", "Consumer brand launch\nData network effects\nTraining data $2M+/mo\nSeries A → global", GREEN),
]
for i, (timeline, name, body, color) in enumerate(phases):
    cx = 0.4 + i * 4.3
    card_rect(sl, cx, 1.95, 4.0, 3.3, CARD, color)
    txt(sl, timeline, cx+0.2, 2.05, 3.6, 0.55, size=9, bold=True, color=color)
    txt(sl, name, cx+0.2, 2.65, 3.6, 0.4, size=16, bold=True, color=WHITE)
    txt(sl, body, cx+0.2, 3.1, 3.6, 1.9, size=12, color=LIGHT)
    if i < 2:
        txt(sl, "→", cx+4.05, 3.35, 0.25, 0.5, size=22, color=color, align=PP_ALIGN.CENTER)

# Acquisition mix
card_rect(sl, 0.4, 5.45, 12.53, 1.7, CARD, VIOLET)
txt(sl, "AGENT ACQUISITION MIX", 0.65, 5.55, 4.0, 0.3, size=10, bold=True, color=VIOLET)
channels = [("Campus Ambassadors", "40%", "$10–20 CAC"), ("Agent Referral", "30%", "$15–25 CAC"), ("Paid Social", "20%", "$30–60 CAC"), ("Organic / App Store", "10%", "$5–15 CAC")]
for i, (ch, pct, cac) in enumerate(channels):
    cx = 0.65 + i * 3.1
    txt(sl, ch, cx, 5.95, 2.8, 0.3, size=11, bold=True, color=WHITE)
    txt(sl, f"{pct}  ·  {cac}", cx, 6.3, 2.8, 0.3, size=10, color=MUTED)
txt(sl, "Blended CAC: $22  ·  Commute-mode agent LTV:CAC = 90:1", 0.65, 6.75, 11.0, 0.3, size=11, bold=True, color=CYAN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=PINK)
section_label(sl, "COMPETITIVE LANDSCAPE", 0.55, 0.28)
slide_headline(sl, "We occupy white space no incumbent has ever entered", 0.5, 0.6, size=28)

# 2×2 matrix
card_rect(sl, 0.4, 1.95, 6.0, 4.8, CARD, VIOLET)
# Axes
txt(sl, "HIGH AI QUALITY", 1.5, 2.05, 4.0, 0.3, size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "LOW AI QUALITY", 1.5, 5.85, 4.0, 0.3, size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "ENTERPRISE\nONLY", 0.45, 3.5, 0.9, 0.6, size=8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(sl, "CONSUMER\n+SMB", 5.35, 3.5, 0.9, 0.6, size=8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
# Center lines
divider_line(sl, 3.2, 4.35, 3.2, MUTED)
box(sl, 3.38, 2.05, 0.04, 4.7, MUTED)
# Competitors
txt(sl, "●  Field Agent", 0.8, 4.7, 2.5, 0.35, size=12, color=ROSE)
txt(sl, "●  BeMyEye", 0.8, 5.1, 2.5, 0.35, size=12, color=ROSE)
txt(sl, "●  Trax", 1.2, 3.7, 2.5, 0.35, size=12, color=ORANGE)
# Mosaic
card_rect(sl, 4.4, 2.35, 1.5, 0.55, VIOLET, CYAN)
txt(sl, "★ MOSAIC", 4.42, 2.38, 1.45, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Feature comparison table
headers = ["", "Field Agent", "BeMyEye", "Trax", "MOSAIC"]
rows = [
    ("Consumer tier", "✗", "✗", "✗", "✓"),
    ("AI extraction", "✗", "✗", "✓", "✓"),
    ("<4hr SLA", "✗", "✗", "✗", "✓"),
    ("Route matching", "✗", "✗", "✗", "✓"),
    ("GDPR on-device", "✗", "?", "✗", "✓"),
    ("Self-serve", "✗", "✗", "✗", "✓"),
]
col_w = [2.2, 1.6, 1.6, 1.3, 1.4]
col_x = [6.55, 8.8, 10.45, 11.8, 12.93]
# But cap at slide width
col_x = [6.55, 8.5, 10.1, 11.3, 12.35]
col_w = [1.9, 1.55, 1.15, 1.0, 0.88]

# Header row
card_rect(sl, 6.5, 1.95, 6.45, 0.42, RGBColor(0x2A, 0x1A, 0x50), VIOLET)
for i, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    color = VIOLET if i==4 else MUTED
    txt(sl, hdr, cx, 2.0, cw, 0.32, size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

for r, row in enumerate(rows):
    cy = 2.42 + r * 0.52
    bg = RGBColor(0x22, 0x14, 0x3A) if r % 2 == 0 else CARD
    card_rect(sl, 6.5, cy, 6.45, 0.5, bg, RGBColor(0x2A, 0x1A, 0x50))
    for i, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
        color = CYAN if (i==4 and val=="✓") else (ROSE if val=="✗" else (MUTED if val=="?" else WHITE))
        txt(sl, val, cx, cy+0.07, cw, 0.35, size=11, bold=(i==0 or i==4), color=color, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — TEAM
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=VIOLET)
section_label(sl, "TEAM", 0.55, 0.28)
slide_headline(sl, "The two people who can build this —\nand the AI team they've assembled", 0.5, 0.6, size=28)

# Founder cards
for i, (role, name, bullets, color) in enumerate([
    ("COMMERCIAL FOUNDER / CEO", "[Founder Name]",
     ["[X] years at [CPG brand] — managed $XM intel budget", "Sold $XM in enterprise SaaS — knows CPG procurement", "[Relevant exits / LinkedIn]"],
     VIOLET),
    ("TECHNICAL FOUNDER / CTO", "[Founder Name]",
     ["Shipped consumer app with XM users — mobile at scale", "AI/ML: real-time inference pipelines, computer vision", "[Relevant exits / GitHub]"],
     CYAN),
]):
    cx = 0.4 + i * 6.35
    card_rect(sl, cx, 2.0, 6.0, 3.0, CARD, color)
    txt(sl, role, cx+0.2, 2.1, 5.6, 0.3, size=9, bold=True, color=color)
    txt(sl, name, cx+0.2, 2.45, 5.6, 0.5, size=18, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        txt(sl, f"→ {b}", cx+0.2, 3.0 + j*0.5, 5.6, 0.4, size=11, color=LIGHT)

# AI Team
card_rect(sl, 0.4, 5.2, 12.53, 1.95, CARD, GREEN)
txt(sl, "THE AI TEAM  ·  Combined output of a 12-person team at $30K/month burn", 0.7, 5.3, 9.0, 0.35, size=11, bold=True, color=GREEN)
ai_tools = [
    ("Claude Code + Cursor", "4× engineering velocity"),
    ("Figma AI + v0", "Design without headcount"),
    ("Pilot.ai + Stripe", "Finance automation"),
    ("Intercom AI", "Customer success at scale"),
]
for i, (tool, desc) in enumerate(ai_tools):
    cx = 0.7 + i * 3.1
    txt(sl, tool, cx, 5.75, 2.8, 0.3, size=11, bold=True, color=WHITE)
    txt(sl, desc, cx, 6.1, 2.8, 0.3, size=10, color=MUTED)
txt(sl, "$30K/month burn vs $400K+ for equivalent traditional team", 0.7, 6.6, 11.5, 0.35, size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — FINANCIAL PROJECTIONS
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=CYAN)
section_label(sl, "FINANCIAL PROJECTIONS", 0.55, 0.28)
slide_headline(sl, "$10M ARR in 18 months — three tiers, three engines", 0.5, 0.6, size=28)

# Revenue table
headers_fin = ["Month", "Consumer", "SMB", "Enterprise", "Total MRR", "ARR Run Rate"]
rows_fin = [
    ("Month 3",  "$6.6K",   "$6.2K",   "$22K",   "$34.8K",  "$418K"),
    ("Month 6",  "$24.2K",  "$24.9K",  "$49.5K", "$98.6K",  "$1.18M"),
    ("Month 12", "$88K",    "$107K",   "$138K",  "$355K",   "$4.25M"),
    ("Month 18", "$220K",   "$242K",   "$275K",  "$822K",   "$9.87M ✓"),
]
col_xs = [0.4, 1.9, 3.7, 5.5, 7.5, 9.6]
col_ws = [1.4, 1.7, 1.7, 1.9, 1.9, 2.3]

card_rect(sl, 0.4, 2.0, 11.6, 0.45, RGBColor(0x2A, 0x1A, 0x50), VIOLET)
for i, (h, cx, cw) in enumerate(zip(headers_fin, col_xs, col_ws)):
    txt(sl, h, cx+0.1, 2.05, cw, 0.32, size=9, bold=True, color=VIOLET if i>=4 else MUTED)

for r, row in enumerate(rows_fin):
    cy = 2.5 + r * 0.6
    is_sa = (r == 3)
    bg = RGBColor(0x1A, 0x0D, 0x30) if is_sa else (RGBColor(0x22, 0x14, 0x3A) if r%2==0 else CARD)
    border = CYAN if is_sa else RGBColor(0x2A, 0x1A, 0x50)
    card_rect(sl, 0.4, cy, 11.6, 0.55, bg, border)
    for i, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
        color = CYAN if (is_sa and i>=4) else (WHITE if i==0 else LIGHT)
        bold = is_sa or i==0
        txt(sl, val, cx+0.1, cy+0.1, cw, 0.35, size=11, bold=bold, color=color)

# Key assumptions + milestones
card_rect(sl, 0.4, 5.0, 5.8, 2.15, CARD, VIOLET)
txt(sl, "KEY ASSUMPTIONS", 0.65, 5.1, 5.0, 0.3, size=10, bold=True, color=VIOLET)
assumptions = ["Agent payout: 40% (fixed)", "Blended gross margin: 55% → 65%+",
               "Enterprise churn: 2.5%/mo", "SMB churn: 5%/mo", "Break-even: Month 20 on seed capital"]
for i, a in enumerate(assumptions):
    txt(sl, f"· {a}", 0.65, 5.5 + i*0.3, 5.3, 0.28, size=10, color=MUTED)

card_rect(sl, 6.45, 5.0, 5.55, 2.15, CARD, CYAN)
txt(sl, "TRAINING DATA UPSIDE (NOT IN MODEL)", 6.7, 5.1, 5.0, 0.3, size=10, bold=True, color=CYAN)
training = [("Month 12", "$22K/mo"), ("Month 18", "$85K/mo"), ("Month 24", "$210K/mo"), ("Month 36", "$2.4M/mo @ 95% margin")]
for i, (mo, rev) in enumerate(training):
    txt(sl, mo, 6.7, 5.5 + i*0.3, 2.0, 0.28, size=10, color=MUTED)
    txt(sl, rev, 8.9, 5.5 + i*0.3, 2.8, 0.28, size=10, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — THE ASK
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=ORANGE)
section_label(sl, "THE ASK", 0.55, 0.28)
slide_headline(sl, "We're raising $5M to prove\nthree cities and three tiers", 0.5, 0.6, size=30)

# Left: raise details
card_rect(sl, 0.4, 2.0, 5.5, 2.2, CARD, ORANGE)
txt(sl, "THE RAISE", 0.65, 2.1, 5.0, 0.3, size=10, bold=True, color=ORANGE)
raise_rows = [("Amount", "$5,000,000"), ("Structure", "Series Seed Preferred"), ("Pre-Money", "$20–25M"), ("Runway", "18 months → Series A")]
for i, (label, val) in enumerate(raise_rows):
    txt(sl, label, 0.65, 2.5 + i*0.4, 2.0, 0.35, size=11, color=MUTED)
    txt(sl, val, 2.8, 2.5 + i*0.4, 3.0, 0.35, size=11, bold=True, color=WHITE)

# Right: use of funds
card_rect(sl, 6.1, 2.0, 7.23, 2.2, CARD, VIOLET)
txt(sl, "USE OF FUNDS", 6.35, 2.1, 6.5, 0.3, size=10, bold=True, color=VIOLET)
funds = [("Engineering (3 hires)", "$1.5M", "30%"), ("Agent acquisition (4 cities)", "$1.0M", "20%"),
         ("ML/AI model dev", "$500K", "10%"), ("Enterprise sales + Marketing", "$1.0M", "20%")]
for i, (cat, amt, pct) in enumerate(funds):
    txt(sl, cat, 6.35, 2.5 + i*0.4, 4.0, 0.35, size=11, color=LIGHT)
    txt(sl, amt, 10.45, 2.5 + i*0.4, 1.2, 0.35, size=11, bold=True, color=WHITE)
    txt(sl, pct, 11.7, 2.5 + i*0.4, 0.8, 0.35, size=11, color=VIOLET)

# Milestones
card_rect(sl, 0.4, 4.4, 12.53, 1.7, CARD, GREEN)
txt(sl, "MILESTONES THIS CAPITAL BUYS", 0.65, 4.5, 6.0, 0.3, size=10, bold=True, color=GREEN)
milestones = [
    "✓ 4 cities live: Atlanta, Dallas, London, Amsterdam",
    "✓ 15,000 active agents",
    "✓ 50 enterprise + 600 SMB subscribers",
    "✓ $10M ARR run rate by Month 18",
    "✓ GDPR compliance architecture complete",
    "✓ Series A raise begins Month 16",
]
for i, m in enumerate(milestones):
    col = i % 2
    row = i // 2
    txt(sl, m, 0.65 + col*6.3, 4.9 + row*0.35, 6.0, 0.32, size=10, color=LIGHT)

# Series A callout
card_rect(sl, 0.4, 6.25, 12.53, 0.9, RGBColor(0x1A, 0x0D, 0x30), ORANGE)
txt(sl, "SERIES A TARGET: $100M at $100M+ pre-money  ·  10× ARR multiple  ·  \"The Field Agent for the AI era\"",
    0.65, 6.45, 12.0, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# APPENDIX: EUROPEAN STRATEGY
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
accent_bar(sl, 0.35, 0.25, color=CYAN)
section_label(sl, "APPENDIX: EUROPEAN STRATEGY", 0.55, 0.28)
slide_headline(sl, "EU-first: less competition, more regulation-readiness", 0.5, 0.6, size=26)

cities_eu = [
    ("🇬🇧 LONDON", "Month 2", "Unilever + Reckitt HQ\nEnglish-native\nICO pragmatic regulator", VIOLET),
    ("🇳🇱 AMSTERDAM", "Month 3", "Heineken HQ\n95% English proficiency\nDutch startup ecosystem", CYAN),
    ("🇩🇪 BERLIN", "Month 5", "Henkel + Beiersdorf\nLargest EU economy\nHigh gig density", GREEN),
    ("🇪🇸 BARCELONA", "Month 7", "Nestlé Iberia\nSpain gig density\nRiders Law manageable", ORANGE),
    ("🇫🇷 PARIS", "Month 9", "L'Oréal + Danone\nHighest CPG density\nRequires full DPIA", PINK),
]
for i, (city, timing, desc, color) in enumerate(cities_eu):
    cx = 0.4 + (i % 3) * 4.25 if i < 3 else 2.55 + (i-3) * 4.25
    cy = 2.1 if i < 3 else 4.4
    card_rect(sl, cx, cy, 3.9, 2.0, CARD, color)
    txt(sl, city, cx+0.15, cy+0.12, 3.6, 0.35, size=13, bold=True, color=color)
    txt(sl, timing, cx+0.15, cy+0.5, 3.6, 0.28, size=10, bold=True, color=WHITE)
    txt(sl, desc, cx+0.15, cy+0.85, 3.6, 0.95, size=10, color=LIGHT)

txt(sl, "BeMyEye (only EU competitor): $35M revenue · enterprise-only · no AI extraction · no consumer tier · 15 years building",
    0.4, 6.5, 12.53, 0.35, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────
output = r"C:\Users\user\projects\PR-001_agentic-ventures\projects\pixie\MOSAIC_Pitch_Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
