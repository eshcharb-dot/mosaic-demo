"""
MOSAIC_PITCH_V6.html → PPTX generator
Captures each slide as a full-res PNG via Playwright,
then inserts as full-bleed images into a 16:9 PPTX.

Run: python generate_pitch_v6_pptx.py
Output: MOSAIC_PITCH_V6.pptx
"""
import sys
import time
from pathlib import Path

BASE = Path(__file__).parent
HTML_FILE = BASE / "MOSAIC_PITCH_V6.html"
OUTPUT_PPTX = BASE / "MOSAIC_PITCH_V6.pptx"
SLIDE_DIR = BASE / "_slide_captures_v6"
SLIDE_DIR.mkdir(exist_ok=True)
TOTAL_SLIDES = 15

# 16:9 presentation dimensions (standard widescreen)
PPTX_WIDTH_EMU  = 12192000   # 1920 × 6350 EMU/px
PPTX_HEIGHT_EMU = 6858000    # 1080 × 6350 EMU/px

def check_deps():
    missing = []
    try: from playwright.sync_api import sync_playwright
    except ImportError: missing.append("playwright")
    try: import pptx
    except ImportError: missing.append("python-pptx")
    if missing:
        print(f"Missing: {', '.join(missing)}")
        print(f"Run: pip install {' '.join(missing)}")
        sys.exit(1)

def capture_slides():
    """Capture slide PNGs — reuse existing ones if already captured by PDF script."""
    from playwright.sync_api import sync_playwright
    existing = sorted(SLIDE_DIR.glob("slide_*.png"))
    if len(existing) == TOTAL_SLIDES:
        print(f"  Reusing {TOTAL_SLIDES} existing captures from {SLIDE_DIR.name}/")
        return existing

    url = HTML_FILE.as_uri()
    print(f"Opening: {url}")
    slides = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto(url, wait_until="networkidle")
        time.sleep(3)
        page.evaluate("""
          const dots = document.getElementById('dots');
          if (dots) dots.style.display = 'none';
          const prog = document.getElementById('progress');
          if (prog) prog.style.display = 'none';
          const counter = document.getElementById('counter');
          if (counter) counter.style.display = 'none';
        """)
        for i in range(TOTAL_SLIDES):
            page.evaluate(f"goTo({i})")
            time.sleep(0.6)
            out = SLIDE_DIR / f"slide_{i+1:02d}.png"
            page.screenshot(path=str(out), full_page=False)
            size_kb = out.stat().st_size // 1024
            print(f"  Slide {i+1:02d}/13  {out.name} ({size_kb}KB)")
            slides.append(out)
        browser.close()
    return slides

def build_pptx(slides):
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width  = Emu(PPTX_WIDTH_EMU)
    prs.slide_height = Emu(PPTX_HEIGHT_EMU)

    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            Emu(0), Emu(0),
            width=Emu(PPTX_WIDTH_EMU),
            height=Emu(PPTX_HEIGHT_EMU)
        )
        print(f"  Slide {i+1:02d}/13 inserted")

    prs.save(str(OUTPUT_PPTX))
    size_mb = OUTPUT_PPTX.stat().st_size / (1024 * 1024)
    print(f"\nOK  PPTX saved: {OUTPUT_PPTX.name} ({size_mb:.1f}MB, {len(slides)} slides)")

if __name__ == "__main__":
    print("=== MOSAIC Pitch V6 — PPTX Generator ===\n")
    check_deps()
    print("Step 1: Capturing slides with Playwright...")
    slides = capture_slides()
    print(f"\nStep 2: Building PPTX from {len(slides)} slides...")
    build_pptx(slides)
    print(f"\nDone. Open: {OUTPUT_PPTX}")
