"""
add_images_to_pptx.py
Adds deck-assets images to MOSAIC_PITCH_V3_DECK.pptx matching the HTML version.
Safe to re-run: removes existing pictures from target slides first.
"""

import os
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

ASSETS = "C:/Users/user/projects/PR-001_agentic-ventures/projects/pixie/deck-assets"
PPTX   = "C:/Users/user/projects/PR-001_agentic-ventures/projects/pixie/MOSAIC_PITCH_V3_DECK.pptx"

prs = Presentation(PPTX)
W = prs.slide_width    # 12192000 EMU (~13.33 in)
H = prs.slide_height   # 6858000  EMU (~7.50 in)

PICTURE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE


def remove_pictures(slide):
    """Remove all picture shapes from a slide."""
    sp_tree = slide.shapes._spTree
    pics = [s for s in slide.shapes if s.shape_type == PICTURE_TYPE]
    for p in pics:
        sp_tree.remove(p._element)
    print(f"  Removed {len(pics)} existing picture(s)")


def add_pic(slide, img_path, left, top, width, height, alpha_pct=100, send_to_back=True):
    """
    Add a picture to a slide with optional transparency.
    alpha_pct: 0=transparent, 100=opaque.
    OOXML: <a:alphaModFix amt='38000'/> inside <a:blip> gives 38% opacity.
    """
    pic = slide.shapes.add_picture(img_path, left, top, width, height)

    if alpha_pct < 100:
        blipFill = pic._element.find(qn('p:blipFill'))
        blip = blipFill.find(qn('a:blip'))
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        alpha_fix = etree.SubElement(blip, f'{{{ns_a}}}alphaModFix')
        alpha_fix.set('amt', str(int(alpha_pct * 1000)))

    if send_to_back:
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)

    return pic


# ── Slide 1: Cover — cover-bg.png full bleed 38% opacity ─────────────────────
print("Slide 1 (Cover): cover-bg.png")
slide1 = prs.slides[0]
remove_pictures(slide1)
add_pic(
    slide1,
    os.path.join(ASSETS, "cover-bg.png"),
    left=0, top=0, width=W, height=H,
    alpha_pct=38, send_to_back=True
)
print("  Added: full bleed, 38% opacity")


# ── Slide 2: Problem — problem-shelf.png left panel ───────────────────────────
# Left panel ~52% wide. Image sits mid-left, ~30% from top, ~38% tall.
print("Slide 2 (Problem): problem-shelf.png")
slide2 = prs.slides[1]
remove_pictures(slide2)

panel_w  = int(W * 0.52)
img_left  = int(W * 0.03)
img_top   = int(H * 0.30)
img_width  = int(panel_w * 0.88)
img_height = int(H * 0.38)

pic2 = add_pic(
    slide2,
    os.path.join(ASSETS, "problem-shelf.png"),
    left=img_left, top=img_top, width=img_width, height=img_height,
    alpha_pct=100, send_to_back=False
)
# Push behind text content but keep visible
sp_tree2 = slide2.shapes._spTree
sp_tree2.remove(pic2._element)
sp_tree2.insert(2, pic2._element)
print("  Added: left panel image area, 100% opacity")


# ── Slide 3: Solution — commute-agent.png left 55% full bleed ────────────────
print("Slide 3 (Solution): commute-agent.png")
slide3 = prs.slides[2]
remove_pictures(slide3)

add_pic(
    slide3,
    os.path.join(ASSETS, "commute-agent.png"),
    left=0, top=0, width=int(W * 0.55), height=H,
    alpha_pct=100, send_to_back=True
)
print("  Added: left 55%, full bleed")


# ── Slide 4: Why Now — gig-workers + ai-extraction card overlays ─────────────
# 3-column card grid inside padding. Middle=gig-workers, Right=ai-extraction.
print("Slide 4 (Why Now): gig-workers.png + ai-extraction.png")
slide4 = prs.slides[3]
remove_pictures(slide4)

pad_x  = int(W * 0.046)
inner_w = W - 2 * pad_x
gap     = int(W * 0.013)
card_w  = (inner_w - 2 * gap) // 3
card_h  = int(H * 0.60)
card_y  = int(H * 0.09) + int(H * 0.18)

mid_x   = pad_x + card_w + gap
right_x = pad_x + 2 * (card_w + gap)

add_pic(slide4, os.path.join(ASSETS, "gig-workers.png"),
        left=mid_x, top=card_y, width=card_w, height=card_h,
        alpha_pct=15, send_to_back=True)
add_pic(slide4, os.path.join(ASSETS, "ai-extraction.png"),
        left=right_x, top=card_y, width=card_w, height=card_h,
        alpha_pct=15, send_to_back=True)
print("  Added: gig-workers (middle card, 15%) + ai-extraction (right card, 15%)")


# ── Slide 5: Product — enterprise-dashboard.png left panel 8% opacity ────────
print("Slide 5 (Product): enterprise-dashboard.png")
slide5 = prs.slides[4]
remove_pictures(slide5)

add_pic(
    slide5,
    os.path.join(ASSETS, "enterprise-dashboard.png"),
    left=0, top=0, width=int(W * 0.48), height=H,
    alpha_pct=8, send_to_back=True
)
print("  Added: left panel, 8% opacity")


# ── Slide 7: Business Model — data-flywheel.png right-edge strip 12% ─────────
print("Slide 7 (Biz Model): data-flywheel.png")
slide7 = prs.slides[6]
remove_pictures(slide7)

bar_h   = int(H * 0.18)
bar_y   = H - bar_h - int(H * 0.08)
strip_w = int(W * 0.16)
strip_x = W - strip_w

add_pic(
    slide7,
    os.path.join(ASSETS, "data-flywheel.png"),
    left=strip_x, top=bar_y, width=strip_w, height=bar_h,
    alpha_pct=12, send_to_back=True
)
print("  Added: right-edge strip of economics bar, 12% opacity")


# ── Slide 8: Traction — city-expansion.png right card 12% opacity ────────────
print("Slide 8 (Traction): city-expansion.png")
slide8 = prs.slides[7]
remove_pictures(slide8)

pad_x8  = int(W * 0.046)
inner_w8 = W - 2 * pad_x8
gap8     = int(W * 0.013)
card_w8  = (inner_w8 - gap8) // 2
grid_y   = int(H * 0.52)
grid_h   = int(H * 0.38)
right_x8 = pad_x8 + card_w8 + gap8

add_pic(
    slide8,
    os.path.join(ASSETS, "city-expansion.png"),
    left=right_x8, top=grid_y, width=card_w8, height=grid_h,
    alpha_pct=12, send_to_back=True
)
print("  Added: right card, 12% opacity")


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(PPTX)
print(f"\nSaved -> {PPTX}")
print("Done. 8 images placed across 7 slides.")
