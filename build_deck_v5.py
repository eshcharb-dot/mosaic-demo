"""
Mosaic Seed Deck V5 — python-pptx builder
15 slides, dark design system, VC-ready
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0A, 0x10)
WHITE    = RGBColor(0xEE, 0xEE, 0xF8)
MUTED    = RGBColor(0x78, 0x78, 0xA0)
PURPLE   = RGBColor(0x7C, 0x6D, 0xF5)
TEAL     = RGBColor(0x00, 0xD4, 0xD4)
GREEN    = RGBColor(0x00, 0xE0, 0x96)
YELLOW   = RGBColor(0xFF, 0xC9, 0x47)
DARK_CARD= RGBColor(0x12, 0x12, 0x1E)
CARD2    = RGBColor(0x16, 0x16, 0x28)

# ── Slide dimensions (16:9) ─────────────────────────────────────────────────
W = Cm(33.87)
H = Cm(19.05)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide  = prs.slides.add_slide(layout)
    return slide

def set_bg(slide, color=BG):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a plain rectangle."""
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()   # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    return shape

def txt(slide, text, x, y, w, h,
        font_size=Pt(18), bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT,
        font_name="Calibri", wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb

def multiline_txt(slide, lines, x, y, w, h,
                  font_size=Pt(18), bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, font_name="Calibri",
                  line_spacing=None):
    """lines is list of (text, font_size, bold, color) tuples or plain strings."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if isinstance(line, str):
            t, fs, b, c = line, font_size, bold, color
        else:
            t, fs, b, c = line
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = t
        run.font.size      = fs
        run.font.bold      = b
        run.font.color.rgb = c
        run.font.name      = font_name
    return txb

def accentline(slide, x, y, w, color=PURPLE, thickness=Pt(2.5)):
    """Thin horizontal accent line."""
    shape = slide.shapes.add_shape(1, x, y, w, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def slide_chrome(slide, slide_num, total=15):
    """Add MOSAIC wordmark top-left and slide number bottom-right."""
    # Wordmark
    txt(slide, "MOSAIC", Cm(0.8), Cm(0.4), Cm(4), Cm(0.7),
        font_size=Pt(11), bold=True, color=MUTED)
    # Slide number
    txt(slide, f"{slide_num} / {total}",
        W - Cm(2.5), H - Cm(0.9), Cm(2.2), Cm(0.7),
        font_size=Pt(10), color=MUTED, align=PP_ALIGN.RIGHT)

def pill(slide, text, x, y, w=Cm(3.5), h=Cm(0.65),
         fill=PURPLE, text_color=WHITE, font_size=Pt(11)):
    """Rounded-ish label pill (rectangle with text)."""
    b = box(slide, x, y, w, h, fill_color=fill)
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return b


def rgb_hex(rgb_color):
    """Convert RGBColor to 6-char hex string."""
    return "%.2X%.2X%.2X" % (rgb_color[0], rgb_color[1], rgb_color[2])


def add_table(slide, data, x, y, w, h,
              header_bg=PURPLE, header_fg=WHITE,
              row_bg=DARK_CARD, row_bg_alt=CARD2,
              cell_fg=WHITE, font_size=Pt(13)):
    """Build a styled table from data (list of row lists)."""
    rows = len(data)
    cols = len(data[0])
    tbl  = slide.shapes.add_table(rows, cols, x, y, w, h).table

    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            # font
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name  = "Calibri"
                    run.font.size  = font_size
                    run.font.bold  = (r == 0)
                    run.font.color.rgb = header_fg if r == 0 else cell_fg
            # bg fill
            bg_col = header_bg if r == 0 else (row_bg if r % 2 == 1 else row_bg_alt)
            hex_val = rgb_hex(bg_col)
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = parse_xml(
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{hex_val}"/>'
                f'</a:solidFill>'
            )
            # remove existing fill if any
            for existing in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(existing)
            tcPr.insert(0, solidFill)
    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def slide_01_cover(prs):
    """COVER"""
    sl = blank_slide(prs)
    set_bg(sl)

    # subtle grid lines (decorative)
    for i in range(0, 16):
        b = box(sl, Cm(i*2.2), Cm(0), Cm(0.02), H, fill_color=RGBColor(0x18,0x18,0x28))

    # Purple glow rectangle behind title
    glow = box(sl, Cm(3), Cm(3.5), Cm(20), Cm(5.5), fill_color=RGBColor(0x0E,0x0C,0x22))
    glow.line.color.rgb = PURPLE
    glow.line.width = Pt(0.5)

    # MOSAIC — giant
    txt(sl, "MOSAIC", Cm(3.5), Cm(3.8), Cm(19), Cm(4.2),
        font_size=Pt(96), bold=True, color=PURPLE,
        align=PP_ALIGN.LEFT)

    # Tagline
    txt(sl, "The physical world, finally legible.",
        Cm(3.5), Cm(8.4), Cm(22), Cm(1.2),
        font_size=Pt(26), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT)

    # Accent line under tagline
    accentline(sl, Cm(3.5), Cm(9.9), Cm(14), color=TEAL, thickness=Pt(2))

    # Subtitle
    txt(sl, "Physical World Intelligence  ·  Seed Round 2026",
        Cm(3.5), Cm(10.3), Cm(22), Cm(0.9),
        font_size=Pt(16), bold=False, color=MUTED,
        align=PP_ALIGN.LEFT)

    # Bottom strip
    box(sl, Cm(0), H - Cm(2.2), W, Cm(2.2), fill_color=RGBColor(0x0D,0x0D,0x1A))
    txt(sl, "Confidential — Not for distribution  |  mosaic-intel.com",
        Cm(1.5), H - Cm(1.8), Cm(22), Cm(0.9),
        font_size=Pt(12), color=MUTED)
    txt(sl, "April 2026",
        W - Cm(5), H - Cm(1.8), Cm(4.5), Cm(0.9),
        font_size=Pt(12), color=MUTED, align=PP_ALIGN.RIGHT)

    slide_chrome(sl, 1)
    return sl


def slide_02_problem(prs):
    """THE PROBLEM"""
    sl = blank_slide(prs)
    set_bg(sl)

    # Title area
    accentline(sl, Cm(1.5), Cm(1.7), Cm(6), color=PURPLE)
    txt(sl, "THE PROBLEM", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=PURPLE)
    txt(sl, "6.8 billion sensors. Zero signal.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.5),
        font_size=Pt(36), bold=True, color=WHITE)

    # Big central stat
    box(sl, Cm(10.5), Cm(2.9), Cm(12.5), Cm(2.1), fill_color=DARK_CARD)
    txt(sl, "$35B/year spent on field intelligence.",
        Cm(11), Cm(3.0), Cm(12), Cm(0.9),
        font_size=Pt(16), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    txt(sl, "Data arrives weeks late. Covers <5% of locations.",
        Cm(11), Cm(3.9), Cm(12), Cm(0.7),
        font_size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

    # Three columns
    col_w = Cm(9.5)
    col_y = Cm(5.4)
    col_h = Cm(9.5)
    cols = [
        ("CPG BRANDS", PURPLE,
         "Lost 4–8% of revenue\nto out-of-stocks they\nnever knew about.\n\nDistributor photos:\n30–40% inaccurate.\n\nField visits cover\n<20% of stores per quarter."),
        ("RETAILERS", TEAL,
         "Can't verify compliance\nwithout a 2-week delay\nand $50–$300/visit.\n\nPlanogram deviations\ngo undetected for weeks.\n\nNo real-time shelf data\nfrom their own estate."),
        ("CONSUMERS", GREEN,
         "Can't trust secondhand\ninfo on crowded venues,\npark conditions, or\nlocal availability.\n\nNo on-demand way to\ncheck physical places\nbefore making a trip."),
    ]
    for i, (header, hcol, body) in enumerate(cols):
        cx = Cm(1.5 + i * 10.5)
        box(sl, cx, col_y, col_w, col_h, fill_color=DARK_CARD)
        box(sl, cx, col_y, col_w, Cm(0.08), fill_color=hcol)
        txt(sl, header, cx + Cm(0.4), col_y + Cm(0.3), col_w - Cm(0.8), Cm(0.8),
            font_size=Pt(15), bold=True, color=hcol)
        accentline(sl, cx + Cm(0.4), col_y + Cm(1.2), col_w - Cm(0.8), color=hcol, thickness=Pt(1))
        txt(sl, body, cx + Cm(0.4), col_y + Cm(1.4), col_w - Cm(0.8), col_h - Cm(1.6),
            font_size=Pt(13.5), color=WHITE)

    slide_chrome(sl, 2)
    return sl


def slide_03_insight(prs):
    """THE INSIGHT"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "THE INSIGHT", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=TEAL)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(6), color=TEAL)

    txt(sl, "We Are Blind.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.5),
        font_size=Pt(40), bold=True, color=WHITE)

    # Paradox statement
    box(sl, Cm(1.5), Cm(2.9), Cm(30.5), Cm(2.4), fill_color=DARK_CARD)
    txt(sl, "6.8 billion smartphones on the planet.\n"
            "Every one of them has a camera, GPS, and a processor more powerful than 2010-era servers.",
        Cm(2.0), Cm(3.1), Cm(29.5), Cm(2.0),
        font_size=Pt(17), color=WHITE)

    txt(sl, "Yet the physical world — its shelves, its crowds, its prices, its conditions — remains almost entirely dark.",
        Cm(1.5), Cm(5.7), Cm(29), Cm(1.5),
        font_size=Pt(18), bold=True, color=YELLOW)

    # Two insight boxes
    box(sl, Cm(1.5), Cm(7.5), Cm(14.5), Cm(4.5), fill_color=CARD2)
    box(sl, Cm(1.5), Cm(7.5), Cm(14.5), Cm(0.06), fill_color=PURPLE)
    txt(sl, "THE OLD MAP", Cm(2.0), Cm(7.7), Cm(13.5), Cm(0.8),
        font_size=Pt(13), bold=True, color=PURPLE)
    txt(sl, "GPS told us WHERE.\n\nReal-time mapping (Google Maps, Waze)\nrevolutionised navigation.\n\nBut WHERE is just geometry.\nIt tells you nothing about WHAT is there.",
        Cm(2.0), Cm(8.5), Cm(13.5), Cm(3.2), font_size=Pt(14), color=WHITE)

    box(sl, Cm(17.2), Cm(7.5), Cm(14.5), Cm(4.5), fill_color=CARD2)
    box(sl, Cm(17.2), Cm(7.5), Cm(14.5), Cm(0.06), fill_color=TEAL)
    txt(sl, "THE MOSAIC LAYER", Cm(17.7), Cm(7.7), Cm(13.5), Cm(0.8),
        font_size=Pt(13), bold=True, color=TEAL)
    txt(sl, "Mosaic tells us WHAT.\n\nIs the shelf stocked? What's the price?\nIs the venue crowded? Is the product compliant?\n\nStructured, AI-extracted, verified.\nDelivered in hours, not weeks.",
        Cm(17.7), Cm(8.5), Cm(13.5), Cm(3.2), font_size=Pt(14), color=WHITE)

    # Quote
    box(sl, Cm(1.5), Cm(12.3), Cm(30.5), Cm(1.3), fill_color=RGBColor(0x10,0x0C,0x22))
    txt(sl, '"GPS told us WHERE. Mosaic tells us WHAT."',
        Cm(2.0), Cm(12.4), Cm(30), Cm(1.1),
        font_size=Pt(20), bold=True, italic=True, color=PURPLE, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 3)
    return sl


def slide_04_solution(prs):
    """THE SOLUTION — 5-step flow"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "THE SOLUTION", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=GREEN)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(6), color=GREEN)
    txt(sl, "Ground truth in 4 hours. Not 4 weeks.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(32), bold=True, color=WHITE)

    # 5 steps horizontal
    steps = [
        (PURPLE, "1", "POST TASK",
         "Enterprise posts\n\"Check oat milk compliance\nat 50 London Tesco stores\""),
        (TEAL,   "2", "MATCH",
         "Mosaic routes to\ncollectors already near\nthose stores on commute"),
        (GREEN,  "3", "CAPTURE",
         "Collector spends 45s.\nOn-device AI validates\nimage quality & redacts faces"),
        (YELLOW, "4", "AI EXTRACT",
         "GPT-4V class model\nextracts structured JSON:\ncompliance scores, flags"),
        (WHITE,  "5", "DELIVERED",
         "Enterprise dashboard:\n4-hour SLA.\n\"6 OOS · 3 misplaced · 41 ✓\""),
    ]

    step_w = Cm(5.8)
    step_h = Cm(8.0)
    start_x = Cm(1.3)
    step_y = Cm(2.9)
    gap = Cm(0.5)

    for i, (col, num, title, body) in enumerate(steps):
        sx = start_x + i * (step_w + gap)
        box(sl, sx, step_y, step_w, step_h, fill_color=DARK_CARD)
        box(sl, sx, step_y, step_w, Cm(0.07), fill_color=col)

        # Step number circle substitute (small box)
        box(sl, sx + Cm(0.25), step_y + Cm(0.25), Cm(0.9), Cm(0.9), fill_color=col)
        txt(sl, num, sx + Cm(0.25), step_y + Cm(0.2), Cm(0.9), Cm(0.9),
            font_size=Pt(16), bold=True, color=BG, align=PP_ALIGN.CENTER)

        txt(sl, title, sx + Cm(0.25), step_y + Cm(1.35), step_w - Cm(0.5), Cm(0.8),
            font_size=Pt(13), bold=True, color=col)
        txt(sl, body, sx + Cm(0.25), step_y + Cm(2.2), step_w - Cm(0.5), Cm(5.2),
            font_size=Pt(12.5), color=WHITE)

        # Arrow between steps
        if i < 4:
            ax = sx + step_w + Cm(0.05)
            txt(sl, "→", ax, step_y + Cm(3.5), gap + Cm(0.2), Cm(0.8),
                font_size=Pt(18), bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Magic numbers bar
    box(sl, Cm(1.3), Cm(11.2), W - Cm(2.6), Cm(2.5), fill_color=CARD2)
    magic = [
        ("90s", "to post a task"),
        ("4hrs", "to deliver results"),
        ("45s", "per collection"),
        ("$12–25", "per store check"),
        (">90%", "AI accuracy"),
    ]
    mw = (W - Cm(2.6)) / len(magic)
    for i, (val, label) in enumerate(magic):
        mx = Cm(1.3) + i * mw
        txt(sl, val, mx, Cm(11.35), mw, Cm(1.1),
            font_size=Pt(24), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        txt(sl, label, mx, Cm(12.4), mw, Cm(0.7),
            font_size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom note
    txt(sl, "Field Agent: $30–175/store, 2-week turnaround, no AI, no dashboard  ·  "
            "Mosaic: $12–25/store, 4-hour SLA, AI extraction, live dashboard",
        Cm(1.5), Cm(14.0), W - Cm(3), Cm(0.9),
        font_size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 4)
    return sl


def slide_05_why_now(prs):
    """WHY NOW"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "WHY NOW", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=YELLOW)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(4), color=YELLOW)
    txt(sl, "Three forces converged in 2024–26 that make Mosaic possible — and defensible.",
        Cm(1.5), Cm(1.0), Cm(29), Cm(1.4),
        font_size=Pt(26), bold=True, color=WHITE)

    forces = [
        (PURPLE, "01", "ON-DEVICE AI",
         "MediaPipe + TFLite run face redaction in real-time on any 2020+ smartphone.\n\n"
         "GDPR compliance without a cloud roundtrip. No biometric data ever leaves the device.\n\n"
         "This wasn't possible before 2024. It's our structural GDPR advantage.",
         "Available 2024+"),
        (TEAL, "02", "GIG ECONOMY MATURITY",
         "50M+ gig workers globally. Already comfortable completing smartphone micro-tasks for micro-pay.\n\n"
         "Deliveroo, Uber, DoorDash spent billions training the supply side.\n\n"
         "We harvest the agent infrastructure they built — at zero marginal cost.",
         "50M+ workers ready"),
        (GREEN, "03", "AI EXTRACTION QUALITY",
         "GPT-4V class models extract structured data from photos with >90% accuracy.\n\n"
         "A shelf photo becomes JSON data in 3 seconds.\n\n"
         "What required a trained human analyst now costs $0.20 of API compute.",
         ">90% accuracy · $0.20/task"),
    ]

    card_w = Cm(9.8)
    card_h = Cm(10.0)
    card_y = Cm(2.8)

    for i, (col, num, title, body, badge) in enumerate(forces):
        cx = Cm(1.3) + i * (card_w + Cm(0.55))
        box(sl, cx, card_y, card_w, card_h, fill_color=DARK_CARD)
        box(sl, cx, card_y, card_w, Cm(0.08), fill_color=col)

        txt(sl, num, cx + Cm(0.4), card_y + Cm(0.3), Cm(1.5), Cm(0.9),
            font_size=Pt(32), bold=True, color=col)
        txt(sl, title, cx + Cm(0.4), card_y + Cm(1.4), card_w - Cm(0.8), Cm(0.8),
            font_size=Pt(15), bold=True, color=col)
        accentline(sl, cx + Cm(0.4), card_y + Cm(2.3), card_w - Cm(0.8), color=col, thickness=Pt(1))
        txt(sl, body, cx + Cm(0.4), card_y + Cm(2.5), card_w - Cm(0.8), card_h - Cm(3.5),
            font_size=Pt(13), color=WHITE)

        # Badge
        box(sl, cx + Cm(0.4), card_y + card_h - Cm(1.2), card_w - Cm(0.8), Cm(0.8),
            fill_color=col)
        txt(sl, badge, cx + Cm(0.4), card_y + card_h - Cm(1.2), card_w - Cm(0.8), Cm(0.8),
            font_size=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)

    # Bottom callout
    box(sl, Cm(1.3), Cm(13.2), W - Cm(2.6), Cm(1.3), fill_color=RGBColor(0x10,0x0A,0x22))
    txt(sl, "The US gig platforms haven't noticed Europe.  ·  BeMyEye has no consumer tier.  ·  "
            "Enterprise AI buyers signed up faster in 2025 than any prior year.",
        Cm(2), Cm(13.3), W - Cm(4), Cm(1.1),
        font_size=Pt(13), color=YELLOW, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 5)
    return sl


def slide_06_collector(prs):
    """THE COLLECTOR — She Was Already There"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "THE COLLECTOR", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=TEAL)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(5.5), color=TEAL)
    txt(sl, '"She Was Already There."',
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(34), bold=True, italic=True, color=WHITE)

    # Persona card (left)
    box(sl, Cm(1.3), Cm(2.9), Cm(13), Cm(9.5), fill_color=DARK_CARD)
    box(sl, Cm(1.3), Cm(2.9), Cm(13), Cm(0.07), fill_color=TEAL)

    # Placeholder avatar
    box(sl, Cm(2.0), Cm(3.2), Cm(3.2), Cm(3.2), fill_color=CARD2)
    txt(sl, "👤", Cm(2.0), Cm(3.3), Cm(3.2), Cm(3.0),
        font_size=Pt(40), align=PP_ALIGN.CENTER, color=MUTED)

    txt(sl, "AMARA", Cm(5.5), Cm(3.4), Cm(8), Cm(0.8),
        font_size=Pt(20), bold=True, color=WHITE)
    txt(sl, "Urban professional · Lagos/London\nAge 28 · Daily transit commuter\n\"I earn £14 before I reach the office.\"",
        Cm(5.5), Cm(4.2), Cm(8.5), Cm(2.5), font_size=Pt(13), color=MUTED)

    accentline(sl, Cm(1.8), Cm(6.7), Cm(12), color=TEAL, thickness=Pt(1))

    txt(sl, "Her commute was already happening.\nMosaic turns dead time into income.\nNo detours. No restructured schedule.",
        Cm(2.0), Cm(7.0), Cm(12), Cm(2.5), font_size=Pt(14), color=WHITE)

    txt(sl, "3 tasks · 45 min commute · £14.50 earned",
        Cm(2.0), Cm(9.8), Cm(12), Cm(0.9),
        font_size=Pt(14), bold=True, color=TEAL)

    # Earnings table (right column)
    box(sl, Cm(15.2), Cm(2.9), Cm(17), Cm(9.5), fill_color=DARK_CARD)
    box(sl, Cm(15.2), Cm(2.9), Cm(17), Cm(0.07), fill_color=PURPLE)
    txt(sl, "COLLECTOR ECONOMICS", Cm(15.5), Cm(3.1), Cm(16), Cm(0.8),
        font_size=Pt(13), bold=True, color=PURPLE)

    earn_data = [
        ["Tier", "Tasks/Month", "Earnings/Month", "Eff. Rate"],
        ["Tier 1 · Power", "60", "~$504", "$14–30/hr*"],
        ["Tier 2 · Regular", "28", "~$210", "$10–18/hr"],
        ["Tier 3 · Occasional", "5", "~$28", "$8–12/hr"],
    ]
    add_table(sl, earn_data,
              Cm(15.5), Cm(4.1), Cm(16.5), Cm(4.0),
              header_bg=PURPLE, font_size=Pt(12))

    txt(sl, "* Marginal rate — commute was happening anyway. Zero opportunity cost.",
        Cm(15.5), Cm(8.3), Cm(16), Cm(0.7), font_size=Pt(11), color=MUTED)

    # LTV:CAC callout
    box(sl, Cm(15.2), Cm(9.3), Cm(8.0), Cm(2.8), fill_color=CARD2)
    txt(sl, "LTV : CAC", Cm(15.5), Cm(9.5), Cm(7.5), Cm(0.8),
        font_size=Pt(13), bold=True, color=MUTED)
    txt(sl, "54 : 1", Cm(15.5), Cm(10.2), Cm(7.5), Cm(1.5),
        font_size=Pt(36), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    box(sl, Cm(24.0), Cm(9.3), Cm(8.2), Cm(2.8), fill_color=CARD2)
    txt(sl, "Blended CAC", Cm(24.3), Cm(9.5), Cm(7.5), Cm(0.8),
        font_size=Pt(13), bold=True, color=MUTED)
    txt(sl, "$22–30", Cm(24.3), Cm(10.2), Cm(7.5), Cm(1.5),
        font_size=Pt(36), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # Acquisition channels
    txt(sl, "Acquisition: Campus ambassadors 40% · Referral 30% · Paid social 20% · Organic 10%",
        Cm(1.5), Cm(13.4), W - Cm(3), Cm(0.8),
        font_size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 6)
    return sl


def slide_07_product(prs):
    """PRODUCT"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "PRODUCT", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=PURPLE)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(4), color=PURPLE)
    txt(sl, "Two apps. One platform. Ground truth on demand.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(28), bold=True, color=WHITE)

    # Collector App panel
    box(sl, Cm(1.3), Cm(2.9), Cm(15), Cm(10.5), fill_color=DARK_CARD)
    box(sl, Cm(1.3), Cm(2.9), Cm(15), Cm(0.07), fill_color=TEAL)
    txt(sl, "COLLECTOR APP — Commute Mode",
        Cm(1.7), Cm(3.1), Cm(14.2), Cm(0.8),
        font_size=Pt(14), bold=True, color=TEAL)

    steps_collect = [
        "Turn on Commute Mode before leaving home",
        "Route analysis pre-loads 3 tasks along your commute",
        "Arrive at location — complete task in 45 seconds",
        "On-device AI validates quality before upload",
        "Earn £14.50 before you reach the office",
    ]
    for i, s in enumerate(steps_collect):
        box(sl, Cm(1.7), Cm(4.1) + i * Cm(1.55), Cm(0.7), Cm(0.7),
            fill_color=TEAL)
        txt(sl, str(i+1), Cm(1.7), Cm(4.05) + i * Cm(1.55), Cm(0.7), Cm(0.7),
            font_size=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)
        txt(sl, s, Cm(2.6), Cm(4.1) + i * Cm(1.55), Cm(13.5), Cm(0.9),
            font_size=Pt(13.5), color=WHITE)

    txt(sl, "Route-matching · Pre-acceptance · Face redaction on-device",
        Cm(1.7), Cm(12.0), Cm(14), Cm(0.8),
        font_size=Pt(12), color=MUTED)

    # Enterprise Dashboard panel
    box(sl, Cm(17.5), Cm(2.9), Cm(15), Cm(10.5), fill_color=DARK_CARD)
    box(sl, Cm(17.5), Cm(2.9), Cm(15), Cm(0.07), fill_color=PURPLE)
    txt(sl, "ENTERPRISE DASHBOARD",
        Cm(17.9), Cm(3.1), Cm(14.2), Cm(0.8),
        font_size=Pt(14), bold=True, color=PURPLE)

    steps_ent = [
        "Post task: location + brief + SLA (90 seconds)",
        "Monitor fulfillment in real-time on map view",
        "AI extraction delivers structured JSON report",
        "Compliance scores, OOS flags, price alerts",
        "Export to your data pipeline via API or CSV",
    ]
    for i, s in enumerate(steps_ent):
        box(sl, Cm(17.9), Cm(4.1) + i * Cm(1.55), Cm(0.7), Cm(0.7),
            fill_color=PURPLE)
        txt(sl, str(i+1), Cm(17.9), Cm(4.05) + i * Cm(1.55), Cm(0.7), Cm(0.7),
            font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, s, Cm(18.8), Cm(4.1) + i * Cm(1.55), Cm(13.5), Cm(0.9),
            font_size=Pt(13.5), color=WHITE)

    txt(sl, "API integration · Webhook alerts · Custom reporting",
        Cm(17.9), Cm(12.0), Cm(14), Cm(0.8),
        font_size=Pt(12), color=MUTED)

    # Magic numbers
    box(sl, Cm(1.3), Cm(13.8), W - Cm(2.6), Cm(1.8), fill_color=CARD2)
    magic = [("90s", "to post"), ("4hrs", "to deliver"), ("45s", "per task"), ("£14.50", "per commute")]
    mw = (W - Cm(2.6)) / len(magic)
    for i, (v, l) in enumerate(magic):
        mx = Cm(1.3) + i * mw
        txt(sl, v, mx, Cm(13.9), mw, Cm(0.9),
            font_size=Pt(22), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        txt(sl, l, mx, Cm(14.8), mw, Cm(0.6),
            font_size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 7)
    return sl


def slide_08_use_cases(prs):
    """TOP 3 USE CASES"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "TOP 3 USE CASES", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=YELLOW)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(6), color=YELLOW)
    txt(sl, "One buyer. Three products. $100K–$400K combined ACV.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(28), bold=True, color=WHITE)

    cases = [
        (PURPLE, "SHELF COMPLIANCE AUDITS",
         "CPG Brand Manager / Field Sales",
         "Covers <20% of stores quarterly.\nDistributor photos 30–40% inaccurate.\nMissed end-cap = $400K–$1.2M lost lift.",
         "50 stores posted → 4hr report → compliance score per store + AI-flagged deviations.",
         "$60K ACV · 60% pilot→contract"),
        (TEAL, "OUT-OF-STOCK DETECTION",
         "Category Manager / Director of Sales Ops",
         "OOS costs brands 4–8% of retail revenue.\nEvents detected 5–10 days after they start.\nPOS data lag: 3–7 days.",
         "Weekly 80-store automated scan. Triggered surge on velocity anomaly. Report by noon Monday.",
         "$70K ACV · 70% pilot→contract"),
        (GREEN, "COMPETITIVE PRICING SCANS",
         "VP Pricing / Revenue Management",
         "Manual audits: 300 stores, 2 weeks, Excel.\nCompetitor promotions run undetected.\nMAP violations undocumented.",
         "Weekly 150-store deploy. AI OCRs shelf tags. Price gap alerts. MAP violation evidence.",
         "$250–300K ACV · 65% pilot→contract"),
    ]

    card_w = Cm(9.8)
    card_h = Cm(10.5)
    for i, (col, title, buyer, pain, solution, metric) in enumerate(cases):
        cx = Cm(1.3) + i * (card_w + Cm(0.55))
        cy = Cm(2.8)
        box(sl, cx, cy, card_w, card_h, fill_color=DARK_CARD)
        box(sl, cx, cy, card_w, Cm(0.07), fill_color=col)

        txt(sl, title, cx + Cm(0.4), cy + Cm(0.2), card_w - Cm(0.8), Cm(0.85),
            font_size=Pt(13), bold=True, color=col)
        txt(sl, f"BUYER: {buyer}", cx + Cm(0.4), cy + Cm(1.15), card_w - Cm(0.8), Cm(0.65),
            font_size=Pt(11), color=MUTED)
        accentline(sl, cx + Cm(0.4), cy + Cm(1.9), card_w - Cm(0.8), color=col, thickness=Pt(1))

        txt(sl, "PAIN", cx + Cm(0.4), cy + Cm(2.05), Cm(1.5), Cm(0.6),
            font_size=Pt(10), bold=True, color=col)
        txt(sl, pain, cx + Cm(0.4), cy + Cm(2.65), card_w - Cm(0.8), Cm(2.5),
            font_size=Pt(12), color=WHITE)

        txt(sl, "MOSAIC", cx + Cm(0.4), cy + Cm(5.3), Cm(2), Cm(0.6),
            font_size=Pt(10), bold=True, color=col)
        txt(sl, solution, cx + Cm(0.4), cy + Cm(5.9), card_w - Cm(0.8), Cm(2.5),
            font_size=Pt(12), color=WHITE)

        # Metric badge
        box(sl, cx + Cm(0.4), cy + card_h - Cm(1.2), card_w - Cm(0.8), Cm(0.9),
            fill_color=col)
        txt(sl, metric, cx + Cm(0.4), cy + card_h - Cm(1.2), card_w - Cm(0.8), Cm(0.9),
            font_size=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)

    txt(sl, "Beachhead: Shelf compliance. Upsell at Month 3: OOS monitoring. Expand at Month 6: Pricing intelligence.",
        Cm(1.5), Cm(13.6), W - Cm(3), Cm(0.9),
        font_size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 8)
    return sl


def slide_09_market(prs):
    """MARKET SIZE"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "MARKET SIZE", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=PURPLE)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(5), color=PURPLE)
    txt(sl, "$35B market. 22% CAGR. Consumer/SMB tier never served.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(26), bold=True, color=WHITE)

    # TAM / SAM / SOM boxes
    market = [
        (PURPLE, "TAM", "$35B",
         "Physical world intelligence:\nRetail intelligence · Mystery shopping\nInsurance field services · AI training data\nBrand compliance · Site intelligence"),
        (TEAL, "SAM", "$8.5B",
         "Crowdsourced, on-demand, AI-native\nvisual intelligence — the segment\naddressable by a smartphone-first platform\nwith consumer + SMB + enterprise tiers"),
        (GREEN, "SOM", "$25M ARR",
         "Credible 3-year target:\n120 enterprise × $66K ACV\n+ 2,000 SMB × $3K ACV\n+ 70K consumers + training data\n(0.36% SAM — achievable)"),
    ]

    mw = Cm(9.8)
    mh = Cm(9.5)
    for i, (col, label, val, desc) in enumerate(market):
        mx = Cm(1.3) + i * (mw + Cm(0.55))
        my = Cm(2.9)
        box(sl, mx, my, mw, mh, fill_color=DARK_CARD)
        box(sl, mx, my, mw, Cm(0.07), fill_color=col)

        txt(sl, label, mx + Cm(0.4), my + Cm(0.25), mw - Cm(0.8), Cm(0.65),
            font_size=Pt(12), bold=True, color=col)
        txt(sl, val, mx + Cm(0.4), my + Cm(0.9), mw - Cm(0.8), Cm(1.8),
            font_size=Pt(42), bold=True, color=WHITE)
        accentline(sl, mx + Cm(0.4), my + Cm(2.85), mw - Cm(0.8), color=col, thickness=Pt(1))
        txt(sl, desc, mx + Cm(0.4), my + Cm(3.0), mw - Cm(0.8), mh - Cm(3.2),
            font_size=Pt(13), color=WHITE)

    # CAGR
    box(sl, Cm(1.3), Cm(12.7), Cm(9.8), Cm(2.0), fill_color=CARD2)
    txt(sl, "CAGR", Cm(1.7), Cm(12.9), Cm(3), Cm(0.7),
        font_size=Pt(13), bold=True, color=MUTED)
    txt(sl, "22%", Cm(1.7), Cm(13.5), Cm(9), Cm(0.9),
        font_size=Pt(28), bold=True, color=YELLOW)

    # Competitor context
    box(sl, Cm(12.2), Cm(12.7), Cm(20), Cm(2.0), fill_color=CARD2)
    txt(sl, "BeMyEye: $35M revenue · enterprise-only · 15 years building  ·  "
            "Field Agent: enterprise-only  ·  Trax: hardware + enterprise only\n"
            "THE CONSUMER/SMB TIER HAS NEVER BEEN SERVED. That is $3–5B of unaddressed market.",
        Cm(12.6), Cm(12.8), Cm(19.2), Cm(1.8),
        font_size=Pt(12.5), color=YELLOW)

    slide_chrome(sl, 9)
    return sl


def slide_10_competitive(prs):
    """COMPETITIVE LANDSCAPE — 2x2 + feature table"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "COMPETITIVE LANDSCAPE", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=TEAL)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(7), color=TEAL)
    txt(sl, "We occupy white space no incumbent has entered.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(28), bold=True, color=WHITE)

    # 2x2 matrix
    mx, my = Cm(1.3), Cm(2.8)
    mw, mh = Cm(14), Cm(11.0)

    # Quadrant fills
    box(sl, mx,        my,         mw/2, mh/2, fill_color=RGBColor(0x10,0x10,0x20))
    box(sl, mx+mw/2,   my,         mw/2, mh/2, fill_color=RGBColor(0x0A,0x0F,0x1A))
    box(sl, mx,        my+mh/2,    mw/2, mh/2, fill_color=RGBColor(0x10,0x10,0x1A))
    box(sl, mx+mw/2,   my+mh/2,    mw/2, mh/2, fill_color=RGBColor(0x08,0x08,0x12))

    # Axes
    box(sl, mx + mw/2 - Cm(0.02), my, Cm(0.04), mh, fill_color=MUTED)
    box(sl, mx, my + mh/2 - Cm(0.02), mw, Cm(0.04), fill_color=MUTED)

    # Axis labels
    txt(sl, "HIGH AI QUALITY", mx, my + Cm(0.1), mw, Cm(0.65),
        font_size=Pt(11), bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txt(sl, "LOW AI QUALITY", mx, my + mh - Cm(0.8), mw, Cm(0.65),
        font_size=Pt(11), bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    txt(sl, "ENTERPRISE\nONLY", mx + Cm(0.2), my + mh/2 - Cm(1.0), Cm(3), Cm(1.0),
        font_size=Pt(10), bold=True, color=MUTED)
    txt(sl, "CONSUMER\n+ SMB", mx + mw - Cm(3.2), my + mh/2 - Cm(1.0), Cm(3), Cm(1.0),
        font_size=Pt(10), bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

    # Competitor dots
    competitors = [
        ("Field Agent", mx + Cm(1.5),   my + mh*0.62, MUTED),
        ("BeMyEye",     mx + Cm(2.5),   my + mh*0.72, MUTED),
        ("Trax",        mx + Cm(1.0),   my + mh*0.55, MUTED),
    ]
    for name, cx, cy, col in competitors:
        box(sl, cx - Cm(0.3), cy - Cm(0.3), Cm(0.6), Cm(0.6), fill_color=col)
        txt(sl, name, cx - Cm(1.5), cy + Cm(0.15), Cm(3), Cm(0.65),
            font_size=Pt(11), color=MUTED)

    # MOSAIC star position — top right quadrant
    box(sl, mx + mw*0.75 - Cm(0.4), my + mh*0.15 - Cm(0.4), Cm(0.8), Cm(0.8),
        fill_color=PURPLE)
    txt(sl, "MOSAIC", mx + mw*0.75 - Cm(1.5), my + mh*0.15 + Cm(0.5), Cm(3.5), Cm(0.7),
        font_size=Pt(13), bold=True, color=PURPLE)

    txt(sl, "White\nspace", mx + mw*0.62, my + mh*0.08, Cm(3.5), Cm(1.0),
        font_size=Pt(11), color=GREEN, italic=True)

    # Feature comparison table (right side)
    feat_data = [
        ["Feature",            "Field Agent", "BeMyEye", "Trax",  "MOSAIC"],
        ["Consumer tier",      "✗",           "✗",       "✗",     "✓"],
        ["AI extraction",      "✗",           "✗",       "partial","✓"],
        ["<4hr SLA",           "✗",           "✗",       "✗",     "✓"],
        ["Route matching",     "✗",           "✗",       "✗",     "✓"],
        ["On-device GDPR",     "✗",           "?",       "✗",     "✓"],
        ["Self-serve",         "✗",           "✗",       "✗",     "✓"],
        ["SMB pricing",        "✗",           "✗",       "✗",     "✓"],
    ]
    add_table(sl, feat_data,
              Cm(16.3), Cm(2.8), Cm(16.5), Cm(9.0),
              header_bg=PURPLE, font_size=Pt(12))

    txt(sl, "No company with strong AI capability has ever built a consumer/SMB product in this space.",
        Cm(1.5), Cm(14.3), W - Cm(3), Cm(0.9),
        font_size=Pt(14), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 10)
    return sl


def slide_11_business_model(prs):
    """BUSINESS MODEL"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "BUSINESS MODEL", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=GREEN)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(6), color=GREEN)
    txt(sl, "Three tiers. One platform. 55%+ gross margin.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(30), bold=True, color=WHITE)

    # Three tier cards
    tiers = [
        (PURPLE, "TIER 1 · CONSUMER",
         "Pay-per-task / Credit wallet",
         "$5–$25 / task",
         "52% gross margin",
         "App store · Viral",
         "Mobile professionals,\nparents, local businesses"),
        (TEAL, "TIER 2 · SMB",
         "$49–$749 / month\nSelf-serve subscription",
         "5–150 tasks/month",
         "62% gross margin",
         "PLG · Content · Shopify app",
         "Indie brands, franchise\noperators, property mgrs"),
        (GREEN, "TIER 3 · ENTERPRISE",
         "$2,500–$50,000+ / month\nAnnual contract",
         "200–10,000+ tasks/month",
         "58% gross margin",
         "Founder sales → AE",
         "CPG brands, retailers,\ninsurers, QSR chains"),
    ]

    tw = Cm(9.8)
    th = Cm(8.5)
    ty = Cm(2.9)

    for i, (col, title, price, volume, margin, acq, who) in enumerate(tiers):
        cx = Cm(1.3) + i * (tw + Cm(0.55))
        box(sl, cx, ty, tw, th, fill_color=DARK_CARD)
        box(sl, cx, ty, tw, Cm(0.07), fill_color=col)

        txt(sl, title, cx + Cm(0.4), ty + Cm(0.2), tw - Cm(0.8), Cm(0.75),
            font_size=Pt(13), bold=True, color=col)
        txt(sl, price, cx + Cm(0.4), ty + Cm(1.05), tw - Cm(0.8), Cm(0.7),
            font_size=Pt(18), bold=True, color=WHITE)
        txt(sl, volume, cx + Cm(0.4), ty + Cm(1.8), tw - Cm(0.8), Cm(0.6),
            font_size=Pt(13), color=MUTED)
        accentline(sl, cx + Cm(0.4), ty + Cm(2.5), tw - Cm(0.8), color=col, thickness=Pt(1))

        box(sl, cx + Cm(0.4), ty + Cm(2.7), tw - Cm(0.8), Cm(0.75), fill_color=col)
        txt(sl, margin, cx + Cm(0.4), ty + Cm(2.7), tw - Cm(0.8), Cm(0.75),
            font_size=Pt(14), bold=True, color=BG, align=PP_ALIGN.CENTER)

        txt(sl, f"ACQ: {acq}", cx + Cm(0.4), ty + Cm(3.6), tw - Cm(0.8), Cm(0.7),
            font_size=Pt(12), color=MUTED)
        txt(sl, who, cx + Cm(0.4), ty + Cm(4.4), tw - Cm(0.8), Cm(1.2),
            font_size=Pt(12), color=WHITE)

    # Revenue table
    rev_data = [
        ["Month",    "MRR",      "Training Data", "Total"],
        ["Mo 3",     "$35K",     "—",             "$35K"],
        ["Mo 6",     "$99K",     "—",             "$99K"],
        ["Mo 12",    "$354K",    "$22K",          "$376K"],
        ["Mo 18",    "$822K",    "$85K",          "$907K"],
        ["Mo 24",    "$1.6M",    "$210K",         "$1.81M"],
    ]
    add_table(sl, rev_data,
              Cm(1.3), Cm(11.7), Cm(20), Cm(4.5),
              header_bg=GREEN, font_size=Pt(13))

    # Flywheel note
    box(sl, Cm(22.3), Cm(11.7), Cm(10.7), Cm(4.5), fill_color=CARD2)
    txt(sl, "TRAINING DATA FLYWHEEL",
        Cm(22.7), Cm(11.9), Cm(10), Cm(0.75),
        font_size=Pt(12), bold=True, color=YELLOW)
    txt(sl, "Every task produces labeled\nreal-world visual data.\n\n"
            "At Mo 12: 22K tasks/day →\nlicenseable training images.\n\n"
            "95% gross margin.\nIncumbents don't have this\n— no consumer tier.",
        Cm(22.7), Cm(12.7), Cm(10), Cm(3.3),
        font_size=Pt(12), color=WHITE)

    # Payout note
    txt(sl, "Agent payout: 40% of task price (fixed)  ·  Platform take: 60% gross  →  AI/infra ~6%  →  Net: 54% blended",
        Cm(1.5), Cm(16.5), W - Cm(3), Cm(0.8),
        font_size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 11)
    return sl


def slide_12_traction(prs):
    """TRACTION"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "TRACTION", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=YELLOW)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(4), color=YELLOW)
    txt(sl, "8 weeks. 2 people. AI-native velocity.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(30), bold=True, color=WHITE)

    # Placeholder metric boxes
    metrics = [
        (PURPLE, "[X]\nVerified Agents", "in [# cities]"),
        (TEAL,   "[X]\nTasks Completed", "in [# weeks]"),
        (GREEN,  "[X]\nPaying Customers", "[X] pilots active"),
        (YELLOW, "$[X]K\nMRR", "Month [X]"),
    ]
    mw = Cm(7.4)
    mh = Cm(4.0)
    for i, (col, val, sub) in enumerate(metrics):
        mx = Cm(1.3) + i * (mw + Cm(0.45))
        box(sl, mx, Cm(2.9), mw, mh, fill_color=DARK_CARD)
        box(sl, mx, Cm(2.9), mw, Cm(0.07), fill_color=col)
        txt(sl, val, mx + Cm(0.3), Cm(3.1), mw - Cm(0.6), Cm(2.2),
            font_size=Pt(28), bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, sub, mx + Cm(0.3), Cm(5.2), mw - Cm(0.6), Cm(0.7),
            font_size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

    # What we built
    box(sl, Cm(1.3), Cm(7.2), Cm(18.5), Cm(7.2), fill_color=DARK_CARD)
    box(sl, Cm(1.3), Cm(7.2), Cm(18.5), Cm(0.07), fill_color=YELLOW)
    txt(sl, "WHAT 2 PEOPLE BUILT IN 8 WEEKS",
        Cm(1.7), Cm(7.4), Cm(18), Cm(0.75),
        font_size=Pt(13), bold=True, color=YELLOW)

    built = [
        "→  Agent app (iOS + Android) — Expo / React Native",
        "→  Customer portal — Next.js + Supabase",
        "→  AI extraction pipeline — GPT-4V + custom validation",
        "→  Route-matching algorithm — geospatial corridor matching",
        "→  3-tier pricing engine + payment processing",
        "→  On-device face redaction — MediaPipe + TFLite",
    ]
    for i, b in enumerate(built):
        txt(sl, b, Cm(1.7), Cm(8.2) + i * Cm(0.95), Cm(18), Cm(0.85),
            font_size=Pt(13.5), color=WHITE)

    # Burn rate
    box(sl, Cm(20.8), Cm(7.2), Cm(11.7), Cm(3.5), fill_color=CARD2)
    txt(sl, "BURN RATE", Cm(21.2), Cm(7.4), Cm(11), Cm(0.7),
        font_size=Pt(13), bold=True, color=MUTED)
    txt(sl, "$30K/month", Cm(21.2), Cm(8.1), Cm(11), Cm(1.2),
        font_size=Pt(32), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, "2-person AI-native team\n(replacing $80K+ in headcount)", Cm(21.2), Cm(9.3), Cm(11), Cm(1.1),
        font_size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

    box(sl, Cm(20.8), Cm(11.1), Cm(11.7), Cm(3.3), fill_color=CARD2)
    txt(sl, "KEY PROOF POINTS", Cm(21.2), Cm(11.3), Cm(11), Cm(0.7),
        font_size=Pt(13), bold=True, color=MUTED)
    proof = [
        "✓  [First customer signed at $X/month]",
        "✓  [X]% task fulfillment rate",
        "✓  [X]% agent retention M1→M2",
        "✓  [Agent NPS / Customer NPS]",
    ]
    for i, p in enumerate(proof):
        txt(sl, p, Cm(21.2), Cm(12.1) + i * Cm(0.55), Cm(11), Cm(0.5),
            font_size=Pt(12), color=WHITE)

    slide_chrome(sl, 12)
    return sl


def slide_13_team(prs):
    """TEAM"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "TEAM", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=PURPLE)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(3), color=PURPLE)
    txt(sl, "The two people who can build this — and the AI team they've assembled.",
        Cm(1.5), Cm(1.0), Cm(30), Cm(1.4),
        font_size=Pt(24), bold=True, color=WHITE)

    # CEO card
    box(sl, Cm(1.3), Cm(2.9), Cm(14.8), Cm(7.0), fill_color=DARK_CARD)
    box(sl, Cm(1.3), Cm(2.9), Cm(14.8), Cm(0.07), fill_color=PURPLE)
    box(sl, Cm(1.8), Cm(3.2), Cm(2.8), Cm(2.8), fill_color=CARD2)
    txt(sl, "CEO", Cm(1.8), Cm(3.4), Cm(2.8), Cm(2.4),
        font_size=Pt(22), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    txt(sl, "[CEO NAME] — Commercial Founder",
        Cm(4.9), Cm(3.2), Cm(10.8), Cm(0.8),
        font_size=Pt(15), bold=True, color=WHITE)
    ceo_bullets = [
        "→  [X] yrs at [CPG company] — managed $Xm competitive intelligence budget",
        "→  Sold $X in enterprise SaaS; knows CPG procurement patterns",
        "→  [Relevant exits / companies]",
    ]
    for i, b in enumerate(ceo_bullets):
        txt(sl, b, Cm(4.9), Cm(4.1) + i * Cm(0.85), Cm(11), Cm(0.75),
            font_size=Pt(13), color=WHITE)
    txt(sl, "linkedin.com/in/[handle]",
        Cm(4.9), Cm(7.0), Cm(11), Cm(0.65),
        font_size=Pt(12), color=MUTED)

    # CTO card
    box(sl, Cm(17.2), Cm(2.9), Cm(14.8), Cm(7.0), fill_color=DARK_CARD)
    box(sl, Cm(17.2), Cm(2.9), Cm(14.8), Cm(0.07), fill_color=TEAL)
    box(sl, Cm(17.7), Cm(3.2), Cm(2.8), Cm(2.8), fill_color=CARD2)
    txt(sl, "CTO", Cm(17.7), Cm(3.4), Cm(2.8), Cm(2.4),
        font_size=Pt(22), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(sl, "[CTO NAME] — Technical Founder",
        Cm(20.8), Cm(3.2), Cm(10.8), Cm(0.8),
        font_size=Pt(15), bold=True, color=WHITE)
    cto_bullets = [
        "→  Shipped [consumer app with X million users] — mobile at scale",
        "→  [AI/ML background] — real-time inference pipelines",
        "→  [Relevant exits / companies]",
    ]
    for i, b in enumerate(cto_bullets):
        txt(sl, b, Cm(20.8), Cm(4.1) + i * Cm(0.85), Cm(11), Cm(0.75),
            font_size=Pt(13), color=WHITE)
    txt(sl, "github.com/[handle]",
        Cm(20.8), Cm(7.0), Cm(11), Cm(0.65),
        font_size=Pt(12), color=MUTED)

    # AI Team
    box(sl, Cm(1.3), Cm(10.3), W - Cm(2.6), Cm(4.0), fill_color=DARK_CARD)
    box(sl, Cm(1.3), Cm(10.3), W - Cm(2.6), Cm(0.07), fill_color=GREEN)
    txt(sl, "THE AI TEAM — $2,400/month replaces $80,000/month in headcount",
        Cm(1.7), Cm(10.5), W - Cm(3.4), Cm(0.8),
        font_size=Pt(14), bold=True, color=GREEN)

    ai_tools = [
        ("Claude Code + Cursor + Devin", "4× engineering velocity"),
        ("v0 + Figma AI", "Design without headcount"),
        ("Pilot.ai + Stripe", "Finance without headcount"),
        ("Intercom AI", "CS without headcount"),
    ]
    tw = (W - Cm(3.4)) / len(ai_tools)
    for i, (tool, desc) in enumerate(ai_tools):
        tx = Cm(1.7) + i * tw
        txt(sl, tool, tx, Cm(11.4), tw - Cm(0.3), Cm(0.75),
            font_size=Pt(13), bold=True, color=WHITE)
        txt(sl, desc, tx, Cm(12.15), tw - Cm(0.3), Cm(0.65),
            font_size=Pt(12), color=MUTED)

    txt(sl, "Combined: the output of a 12-person team at a $30K/month burn rate.",
        Cm(1.7), Cm(13.1), W - Cm(3.4), Cm(0.7),
        font_size=Pt(13), bold=True, color=YELLOW)

    slide_chrome(sl, 13)
    return sl


def slide_14_financials(prs):
    """FINANCIALS"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "FINANCIALS", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=TEAL)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(5), color=TEAL)
    txt(sl, "Bottom-up projections. Three tiers. Series A trigger at Month 18.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(24), bold=True, color=WHITE)

    # MRR table
    mrr_data = [
        ["Month", "T1 Consumer", "T2 SMB", "T3 Enterprise", "MRR Total", "Training Data", "Grand Total"],
        ["Mo 3",  "$6,600",     "$6,225",  "$22,000",       "$34,825",  "—",             "$34,825"],
        ["Mo 6",  "$24,200",    "$24,900", "$49,500",       "$98,600",  "—",             "$98,600"],
        ["Mo 9",  "$44,000",    "$57,200", "$82,500",       "$183,700", "—",             "$183,700"],
        ["Mo 12", "$88,000",    "$107,070","$137,500",      "$354,570", "$22,000",       "$376,570"],
        ["Mo 15", "$143,000",   "$170,100","$192,500",      "$505,600", "$52,000",       "$557,600"],
        ["Mo 18", "$220,000",   "$242,100","$275,000",      "$822,100", "$85,000",       "$907,100"],
        ["Mo 21", "$385,000",   "$319,000","$385,000",      "$1,089,000","$145,000",     "$1,234,000"],
        ["Mo 24", "$550,000",   "$398,800","$467,500",      "$1,626,300","$210,000",     "$1,836,300"],
    ]
    add_table(sl, mrr_data,
              Cm(1.3), Cm(2.9), W - Cm(2.6), Cm(7.8),
              header_bg=TEAL, font_size=Pt(12))

    # Key metrics below table
    kpis = [
        (PURPLE, "Month 18 ARR", "$10.9M", "Series A trigger ✓"),
        (TEAL,   "Blended Gross Margin", "55% → 65%+", "at scale"),
        (GREEN,  "Burn Rate (Post-Seed)", "$120K/month", "team expands to 8"),
        (YELLOW, "Break-even", "Month 20", "on seed capital"),
    ]
    kw = (W - Cm(2.6)) / len(kpis)
    for i, (col, label, val, sub) in enumerate(kpis):
        kx = Cm(1.3) + i * kw
        box(sl, kx, Cm(11.1), kw - Cm(0.2), Cm(3.4), fill_color=DARK_CARD)
        box(sl, kx, Cm(11.1), kw - Cm(0.2), Cm(0.06), fill_color=col)
        txt(sl, label, kx + Cm(0.3), Cm(11.2), kw - Cm(0.6), Cm(0.7),
            font_size=Pt(12), bold=True, color=col)
        txt(sl, val, kx + Cm(0.3), Cm(11.9), kw - Cm(0.6), Cm(1.3),
            font_size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, sub, kx + Cm(0.3), Cm(13.2), kw - Cm(0.6), Cm(0.65),
            font_size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

    txt(sl, "Key assumptions: 40% agent payout · 2.5% enterprise churn/mo · 5% SMB churn/mo · "
            "3-city launch Mo 4 · 8 cities Mo 12 · MAC at Mo 12: 3,000–5,000",
        Cm(1.5), Cm(14.8), W - Cm(3), Cm(0.9),
        font_size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 14)
    return sl


def slide_15_ask(prs):
    """THE ASK"""
    sl = blank_slide(prs)
    set_bg(sl)

    txt(sl, "THE ASK", Cm(1.5), Cm(0.5), Cm(20), Cm(0.8),
        font_size=Pt(12), bold=True, color=GREEN)
    accentline(sl, Cm(1.5), Cm(1.7), Cm(3.5), color=GREEN)
    txt(sl, "We're raising $5M to prove three cities and three tiers.",
        Cm(1.5), Cm(1.0), Cm(28), Cm(1.4),
        font_size=Pt(28), bold=True, color=WHITE)

    # Round terms
    terms = [
        (PURPLE, "RAISING", "$5,000,000"),
        (TEAL,   "PRE-MONEY", "$20M"),
        (GREEN,  "DILUTION", "~20%"),
        (YELLOW, "RUNWAY", "18 months"),
    ]
    tw = (W - Cm(2.6)) / len(terms)
    for i, (col, label, val) in enumerate(terms):
        tx = Cm(1.3) + i * tw
        box(sl, tx, Cm(2.8), tw - Cm(0.3), Cm(2.4), fill_color=DARK_CARD)
        box(sl, tx, Cm(2.8), tw - Cm(0.3), Cm(0.07), fill_color=col)
        txt(sl, label, tx + Cm(0.3), Cm(2.95), tw - Cm(0.9), Cm(0.65),
            font_size=Pt(12), bold=True, color=col)
        txt(sl, val, tx + Cm(0.3), Cm(3.55), tw - Cm(0.9), Cm(1.2),
            font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Use of funds
    box(sl, Cm(1.3), Cm(5.6), Cm(17.0), Cm(9.2), fill_color=DARK_CARD)
    box(sl, Cm(1.3), Cm(5.6), Cm(17.0), Cm(0.07), fill_color=PURPLE)
    txt(sl, "USE OF FUNDS",
        Cm(1.7), Cm(5.75), Cm(16.5), Cm(0.75),
        font_size=Pt(13), bold=True, color=PURPLE)

    funds = [
        ("Engineering team (3 hires)",       "$1,500,000", "30%"),
        ("Agent acquisition (4 cities)",      "$1,000,000", "20%"),
        ("ML / AI model development",         "$500,000",   "10%"),
        ("Enterprise sales (1 AE)",           "$500,000",   "10%"),
        ("Marketing (SMB + consumer)",        "$500,000",   "10%"),
        ("Legal / compliance (GDPR, SOC 2)", "$300,000",   " 6%"),
        ("Operations (city launches)",        "$400,000",   " 8%"),
        ("Buffer",                            "$300,000",   " 6%"),
    ]
    for i, (item, amount, pct) in enumerate(funds):
        fy = Cm(6.65) + i * Cm(0.97)
        txt(sl, item,   Cm(1.7),  fy, Cm(10.5), Cm(0.85), font_size=Pt(13), color=WHITE)
        txt(sl, amount, Cm(12.0), fy, Cm(3.5),  Cm(0.85), font_size=Pt(13), bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
        txt(sl, pct,    Cm(15.7), fy, Cm(2.0),  Cm(0.85), font_size=Pt(13), color=MUTED, align=PP_ALIGN.RIGHT)

    # Milestones
    box(sl, Cm(19.5), Cm(5.6), Cm(13.5), Cm(9.2), fill_color=DARK_CARD)
    box(sl, Cm(19.5), Cm(5.6), Cm(13.5), Cm(0.07), fill_color=GREEN)
    txt(sl, "MILESTONES THIS CAPITAL BUYS",
        Cm(19.9), Cm(5.75), Cm(13), Cm(0.75),
        font_size=Pt(13), bold=True, color=GREEN)

    milestones = [
        "✓  4 cities live (Atlanta, Dallas, London, Amsterdam)",
        "✓  15,000 registered · 3,000–5,000 monthly active collectors",
        "✓  50 enterprise customers + 600 SMB subscribers",
        "✓  $3.5M ARR by Month 12 (stretch: $4.2M)",
        "✓  $10M+ ARR run rate trigger → Series A",
        "✓  GDPR compliance architecture completed",
        "✓  Series A raise begins Month 16",
    ]
    for i, m in enumerate(milestones):
        txt(sl, m, Cm(19.9), Cm(6.55) + i * Cm(1.1), Cm(12.8), Cm(0.95),
            font_size=Pt(13), color=WHITE)

    # Closing statement
    box(sl, Cm(1.3), Cm(15.1), W - Cm(2.6), Cm(1.8), fill_color=RGBColor(0x0C,0x08,0x1E))
    txt(sl, '"The physical world is ready to answer."',
        Cm(2), Cm(15.3), W - Cm(4), Cm(1.3),
        font_size=Pt(26), bold=True, italic=True, color=PURPLE,
        align=PP_ALIGN.CENTER)

    txt(sl, "Series A at Month 18: $10M ARR · 10× revenue multiple → $100M valuation · 3× YoY growth",
        Cm(1.5), Cm(17.1), W - Cm(3), Cm(0.8),
        font_size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

    slide_chrome(sl, 15)
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    output_path = r"C:\Users\user\projects\PR-001_agentic-ventures\projects\pixie\MOSAIC_SEED_DECK_V5.pptx"

    prs = new_prs()

    print("Building slide 1 — Cover")
    slide_01_cover(prs)
    print("Building slide 2 — Problem")
    slide_02_problem(prs)
    print("Building slide 3 — Insight")
    slide_03_insight(prs)
    print("Building slide 4 — Solution")
    slide_04_solution(prs)
    print("Building slide 5 — Why Now")
    slide_05_why_now(prs)
    print("Building slide 6 — Collector")
    slide_06_collector(prs)
    print("Building slide 7 — Product")
    slide_07_product(prs)
    print("Building slide 8 — Use Cases")
    slide_08_use_cases(prs)
    print("Building slide 9 — Market")
    slide_09_market(prs)
    print("Building slide 10 — Competitive")
    slide_10_competitive(prs)
    print("Building slide 11 — Business Model")
    slide_11_business_model(prs)
    print("Building slide 12 — Traction")
    slide_12_traction(prs)
    print("Building slide 13 — Team")
    slide_13_team(prs)
    print("Building slide 14 — Financials")
    slide_14_financials(prs)
    print("Building slide 15 — The Ask")
    slide_15_ask(prs)

    prs.save(output_path)
    print(f"\nSaved -> {output_path}")
    print(f"Slide count: {len(prs.slides)}")

if __name__ == "__main__":
    main()
